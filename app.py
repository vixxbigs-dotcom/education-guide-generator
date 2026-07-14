import base64
import html
import io
import json
import re
import time
from concurrent.futures import ThreadPoolExecutor
from datetime import datetime
from pathlib import Path
from urllib.parse import quote, urlencode
from urllib.request import Request, urlopen
from urllib.error import HTTPError, URLError

import streamlit as st
import streamlit.components.v1 as components


st.set_page_config(page_title="교육 안내문 작성 도구 (Beta)", page_icon="✉️", layout="wide")


def load_local_env_file_once() -> None:
    """python-dotenv 없이 프로젝트 루트의 .env 값을 환경변수로 읽어옵니다."""
    import os
    env_path = Path.cwd() / ".env"
    if not env_path.exists():
        return
    try:
        for raw_line in env_path.read_text(encoding="utf-8").splitlines():
            line = raw_line.strip()
            if not line or line.startswith("#") or "=" not in line:
                continue
            key, value = line.split("=", 1)
            key = key.strip()
            value = value.strip().strip('"').strip("'")
            if key and key not in os.environ:
                os.environ[key] = value
    except Exception:
        pass


load_local_env_file_once()


# -----------------------------
# 기본 입력값 / 초기화 유틸
# -----------------------------
DEFAULT_CURRICULUM = [
    {"day": "Day1", "time": "10:00 ~ 12:00", "subject": "M1. 팀빌딩과 회고", "speaker": ""},
    {"day": "Day1", "time": "13:00 ~ 16:00", "subject": "M2. 이전부터 오늘까지 나", "speaker": ""},
    {"day": "Day1", "time": "16:00 ~ 18:00", "subject": "M3. 내일의 나 (경력 적응성의 관점)", "speaker": ""},
    {"day": "Day1", "time": "18:00 ~ 19:00", "subject": "석식 만찬", "speaker": ""},
    {"day": "Day2", "time": "09:00 ~ 12:00", "subject": "M4. HMS/RESPECT", "speaker": "사내강사"},
    {"day": "Day2", "time": "13:00 ~ 15:00", "subject": "M5. 잡 크래프팅과 일의 의미, 몰입", "speaker": ""},
    {"day": "Day2", "time": "15:00 ~ 17:00", "subject": "M6. 개인의 목표와 조직의 목표 align", "speaker": ""},
]

DEFAULT_CONTACTS = [
    {"role": "현장운영자 000 프로", "phone": "010-0000-0000"},
    {"role": "멀티캠퍼스 000 프로", "phone": "010-0000-0000"},
    {"role": "00그룹 000 책임", "phone": "010-0000-0000"},
]

DEFAULT_VALUES = {
    "company_name": "",
    "course_name": "멀티캠퍼스 입문 교육",
    "delivery_mode": "대면",
    "delivery_custom": "",
    "welcome_title": "입과를 환영합니다!",
    "welcome_body_text": "{교육명}\n강의에 입과하신 여러분 환영합니다! 해당 강의는 {운영방식}으로 진행되며,\n하기 내용을 사전에 꼭 확인하신 후 입과해주시길 부탁드립니다.",
    "time_notice_text": "특히, 교육 시작시간은 1일차 {1일차}시, 2일차 {2일차}시이니 일정 확인 부탁드립니다.",
    "app_theme": "다크 모드",
    "main_color_picker": "#0088C9",
    "footer_color_picker": "#00A651",
    "curr_header_text_color_picker": "#FFFFFF",
    "use_custom_footer_color": False,
    "logo_position": "우측 상단",
    "logo_max_height": 52,
    "font_mode": "assets/fonts 폴더",
    "primary_font_select": "Malgun Gothic",
    "primary_font_custom": "Pretendard",
    "date_range": "6/16(월)~6/17(화)",
    "day1_time": "10",
    "day2_time": "9",
    "place_name": "스타필드 수원 타임체임버",
    "road_address": "",
    "display_place_name": "스타필드 수원 타임체임버",
    "display_road_address": "",
    "display_location_text": "스타필드 수원 타임체임버",
    "last_address_fetch_message": "",
    "selected_map_capture_variant": "기본 크기 (1000×800 / 5:4)",
    "naver_background_message": "",
    "naver_map_background_message": "",
    "naver_address_background_message": "",
    "naver_client_id_input": "",
    "naver_client_secret_input": "",
    "naver_local_results": [],
    "naver_local_search_message": "",
    "selected_naver_local_index": -1,
    "show_curriculum_table": True,
    "curriculum_title": "상세 커리큘럼",
    "curriculum_columns_text": "Day, 시간, 교육 내용, 강사/비고",
    "curriculum_column_defs": [
        {"column_name": "Day"},
        {"column_name": "시간"},
        {"column_name": "교육 내용"},
        {"column_name": "강사/비고"},
    ],
    "info_text": "- 숙소는 2인 1실로 제공될 예정이며 생수, 비누, 샴푸, 헤어드라이기, 냉장고, TV, 전화기, 비데, 유무선인터넷 등이 구비되어 있습니다.\n- 1일차 석식은 연수원 외부에서 진행될 예정이오니 참고 부탁드립니다.",
    "preview_scale_percent": 72,
    "export_file_name": "education_notice_photo_style",
    "captured_map_data_url": "",
    "captured_map_file_path": "",
    "captured_map_files": {},
    "last_map_capture_message": "",
    "pasted_map_data_url": "",
    "pending_location_update": None,
    "zoom_url": "",
    "overview_section_title": "교육 개요",
    "location_section_title": "교육 장소",
    "notice_section_title": "안내 사항",
    "contact_section_title": "관련 문의",
    "curriculum_odd_row_color_picker": "#FFFFFF",
    "curriculum_even_row_color_picker": "#FFFFFF",
    "excel_import_message": "",
}

WIDGET_KEYS_TO_CLEAR_ON_RESET = [
    "logo_file_uploader",
    "map_file_uploader",
    "curr_editor",
    "curriculum_column_defs_editor",
    "contacts_editor",
    "font_folder_select",
    "curriculum_excel_uploader",
]


def init_defaults() -> None:
    for key, value in DEFAULT_VALUES.items():
        if key not in st.session_state:
            st.session_state[key] = value
    if "curriculum" not in st.session_state:
        st.session_state.curriculum = [row.copy() for row in DEFAULT_CURRICULUM]
    if "contacts" not in st.session_state:
        st.session_state.contacts = [row.copy() for row in DEFAULT_CONTACTS]


def order_curriculum_columns(columns: list[str]) -> list[str]:
    """시간표 기본 열 순서를 Day → 시간 → 나머지 순서로 정리합니다."""
    cleaned: list[str] = []
    for col in columns:
        clean = str(col or "").strip()
        if clean == "일차":
            clean = "Day"
        if clean and clean not in cleaned:
            cleaned.append(clean)

    ordered: list[str] = []
    for priority in ["Day", "시간"]:
        if priority in cleaned:
            ordered.append(priority)
    ordered.extend([col for col in cleaned if col not in ordered])
    return ordered or ["Day", "시간", "교육 내용", "강사/비고"]


def migrate_curriculum_day_defaults() -> None:
    """기존 세션에 남아 있는 '일차/1일차/2일차' 기본값과 열 순서를 Day 표기로 정리합니다."""
    column_defs = st.session_state.get("curriculum_column_defs", [])
    raw_columns = []
    changed_defs = False
    for row in to_records(column_defs):
        col = str(row.get("column_name", "") or "").strip()
        if col == "일차":
            col = "Day"
            changed_defs = True
        if col:
            raw_columns.append(col)

    ordered_columns = order_curriculum_columns(raw_columns)
    if raw_columns and ordered_columns != raw_columns:
        changed_defs = True
    if ordered_columns and changed_defs:
        migrated_defs = [{"column_name": col} for col in ordered_columns]
        st.session_state.curriculum_column_defs = migrated_defs
        st.session_state.curriculum_columns_text = ", ".join(ordered_columns)

    curriculum_rows = []
    changed_rows = False
    for row in to_records(st.session_state.get("curriculum", [])):
        new_row = dict(row)
        for key in ["day", "Day", "일차"]:
            value = str(new_row.get(key, "") or "").strip()
            if value == "1일차":
                new_row[key] = "Day1"
                changed_rows = True
            elif value == "2일차":
                new_row[key] = "Day2"
                changed_rows = True
        curriculum_rows.append(new_row)
    if curriculum_rows and changed_rows:
        st.session_state.curriculum = curriculum_rows


def reset_all_fields() -> None:
    current_app_theme = st.session_state.get("app_theme", DEFAULT_VALUES["app_theme"])
    for key, value in DEFAULT_VALUES.items():
        st.session_state[key] = value
    st.session_state.app_theme = current_app_theme
    st.session_state.curriculum = [row.copy() for row in DEFAULT_CURRICULUM]
    st.session_state.contacts = [row.copy() for row in DEFAULT_CONTACTS]
    for key in WIDGET_KEYS_TO_CLEAR_ON_RESET:
        st.session_state.pop(key, None)
    st.session_state.pop("last_captured_map_file", None)
    st.session_state.pop("captured_map_file_path", None)
    st.session_state.pop("captured_map_files", None)
    st.session_state.pop("last_map_capture_message", None)
    st.session_state.pop("naver_combined_future", None)
    st.session_state.pop("naver_combined_started_at", None)
    st.session_state.pop("naver_map_future", None)
    st.session_state.pop("naver_map_started_at", None)
    st.session_state.pop("naver_address_future", None)
    st.session_state.pop("naver_address_started_at", None)
    for key in list(st.session_state.keys()):
        if key.startswith("white_grid_") or key.endswith("_white_rows"):
            st.session_state.pop(key, None)
    st.cache_data.clear()
    st.cache_resource.clear()



# -----------------------------
# 기본 유틸
# -----------------------------
def esc(value: object) -> str:
    return html.escape(str(value or ""), quote=False)


def esc_attr(value: object) -> str:
    return html.escape(str(value or ""), quote=True)


def normalize_bullet(line: str) -> str:
    line = str(line or "").strip()
    if line.startswith("-"):
        return line[1:].strip()
    return line


def to_records(data: object) -> list[dict]:
    """st.data_editor 반환값이 list 또는 DataFrame이어도 dict 리스트로 통일합니다."""
    if data is None:
        return []
    if hasattr(data, "to_dict"):
        try:
            records = data.to_dict("records")
            return [row for row in records if isinstance(row, dict)]
        except TypeError:
            pass
    if isinstance(data, list):
        return [row for row in data if isinstance(row, dict)]
    return []


def file_to_data_url(uploaded_file) -> str:
    if uploaded_file is None:
        return ""
    file_bytes = uploaded_file.getvalue()
    mime_type = uploaded_file.type or "image/png"
    return f"data:{mime_type};base64,{base64.b64encode(file_bytes).decode('ascii')}"


def get_captures_dir() -> Path:
    return Path(__file__).resolve().parent / "assets" / "captures"


def image_bytes_to_data_url(image_bytes: bytes, mime_type: str = "image/png") -> str:
    return f"data:{mime_type};base64,{base64.b64encode(image_bytes).decode('ascii')}"


def file_path_to_data_url(file_path: str) -> str:
    path = Path(str(file_path or "")).expanduser()
    if not file_path or not path.exists() or not path.is_file():
        return ""
    suffix = path.suffix.lower()
    mime_type = "image/png"
    if suffix in [".jpg", ".jpeg"]:
        mime_type = "image/jpeg"
    elif suffix == ".gif":
        mime_type = "image/gif"
    elif suffix == ".webp":
        mime_type = "image/webp"
    return f"data:{mime_type};base64,{base64.b64encode(path.read_bytes()).decode('ascii')}"


def normalize_image_data_url(value: str) -> str:
    """클립보드 붙여넣기로 들어온 이미지 data URL만 안내문 이미지로 사용합니다."""
    clean = str(value or "").strip()
    if clean.startswith("data:image/") and ";base64," in clean:
        return clean
    return ""



def strip_html_tags(value: str) -> str:
    """네이버 지역 검색 API title에 섞이는 <b> 태그 등을 제거합니다."""
    text = re.sub(r"<[^>]+>", "", str(value or ""))
    return html.unescape(text).strip()


def normalize_naver_coord(value: object) -> float | None:
    """네이버 지역 검색 API의 mapx/mapy 정수 좌표를 경도/위도 실수로 변환합니다."""
    try:
        number = float(str(value).strip())
    except (TypeError, ValueError):
        return None
    # 지역 검색 API는 WGS84 기준 정수형 좌표를 반환합니다. 보통 1e7로 나누면 경도/위도가 됩니다.
    if abs(number) > 1000:
        number = number / 10_000_000
    return number


def get_naver_local_api_credentials() -> tuple[str, str]:
    """.env/환경변수 > st.secrets 순서로 네이버 검색 API 키를 가져옵니다.

    최종 사용자는 API 키를 화면에 입력하지 않도록, 기본값은 프로젝트 루트의
    .env 파일 또는 실행 환경변수에서 읽습니다.
    """
    try:
        import os
        env_id = os.getenv("NAVER_CLIENT_ID", "").strip()
        env_secret = os.getenv("NAVER_CLIENT_SECRET", "").strip()
        if env_id and env_secret:
            return env_id, env_secret
    except Exception:
        pass

    try:
        secret_id = str(st.secrets.get("NAVER_CLIENT_ID", "") or "").strip()
        secret_secret = str(st.secrets.get("NAVER_CLIENT_SECRET", "") or "").strip()
        if secret_id and secret_secret:
            return secret_id, secret_secret
    except Exception:
        pass

    return "", ""


def has_naver_local_api_credentials() -> bool:
    client_id, client_secret = get_naver_local_api_credentials()
    return bool(client_id and client_secret)


def queue_location_update(title: str = "", address: str = "", message: str = "") -> None:
    """위젯 생성 이후에는 key 값을 직접 바꾸지 않고, 다음 rerun 초기에 반영하도록 임시 저장합니다."""
    st.session_state.pending_location_update = {
        "title": str(title or "").strip(),
        "address": str(address or "").strip(),
        "message": str(message or "").strip(),
    }


def apply_pending_location_update() -> None:
    pending = st.session_state.get("pending_location_update")
    if not pending:
        return
    title = str((pending or {}).get("title", "") or "").strip()
    address = str((pending or {}).get("address", "") or "").strip()
    message = str((pending or {}).get("message", "") or "").strip()
    if title:
        st.session_state.place_name = title
    if address:
        st.session_state.road_address = address
    if title or address:
        st.session_state.display_location_text = format_location_line(title, address)
    if message:
        st.session_state.last_address_fetch_message = message
    st.session_state.pending_location_update = None

@st.cache_data(show_spinner=False, ttl=3600)
def search_naver_local_api(query: str, client_id: str, client_secret: str, display: int = 5) -> dict:
    """네이버 Developers 검색 > 지역 API로 장소 후보를 가져옵니다."""
    query = str(query or "").strip()
    if not query:
        return {"ok": False, "message": "검색어를 입력해 주세요.", "items": []}
    if not client_id or not client_secret:
        return {"ok": False, "message": "장소 검색 설정을 확인해 주세요.", "items": []}

    params = urlencode({"query": query, "display": max(1, min(int(display or 5), 5)), "start": 1, "sort": "random"})
    api_url = f"https://openapi.naver.com/v1/search/local.json?{params}"
    req = Request(
        api_url,
        headers={
            "X-Naver-Client-Id": client_id,
            "X-Naver-Client-Secret": client_secret,
            "User-Agent": "education-notice-generator/1.0",
        },
        method="GET",
    )

    try:
        with urlopen(req, timeout=8) as response:
            payload = json.loads(response.read().decode("utf-8"))
    except HTTPError as exc:
        detail = ""
        try:
            detail = exc.read().decode("utf-8")[:500]
        except Exception:
            detail = ""
        return {"ok": False, "message": f"네이버 지역 검색 API 오류: HTTP {exc.code} {detail}", "items": []}
    except URLError as exc:
        return {"ok": False, "message": f"네이버 지역 검색 API 연결 실패: {exc}", "items": []}
    except Exception as exc:
        return {"ok": False, "message": f"네이버 지역 검색 API 처리 실패: {exc}", "items": []}

    items = []
    for raw in payload.get("items", []) or []:
        title = strip_html_tags(raw.get("title", ""))
        road = strip_html_tags(raw.get("roadAddress", ""))
        addr = strip_html_tags(raw.get("address", ""))
        category = strip_html_tags(raw.get("category", ""))
        link = str(raw.get("link", "") or "").strip()
        lon = normalize_naver_coord(raw.get("mapx"))
        lat = normalize_naver_coord(raw.get("mapy"))
        items.append({
            "title": title,
            "category": category,
            "roadAddress": road,
            "address": addr,
            "mapx": raw.get("mapx", ""),
            "mapy": raw.get("mapy", ""),
            "lon": lon,
            "lat": lat,
            "link": link,
        })
    return {"ok": True, "message": f"검색 결과 {len(items)}건을 가져왔습니다.", "items": items}


def apply_naver_local_item(item: dict) -> None:
    """검색 결과 카드에서 선택한 장소를 다음 rerun 때 안전하게 반영합니다."""
    title = strip_html_tags(item.get("title", ""))
    road = strip_html_tags(item.get("roadAddress", ""))
    addr = strip_html_tags(item.get("address", ""))
    chosen_address = road or addr
    queue_location_update(title, chosen_address, "선택한 장소를 안내문에 반영했습니다.")
    st.session_state.naver_local_search_message = ""


def apply_first_naver_local_result() -> None:
    """검색어 기준 첫 번째 장소 후보를 안내문 장소/주소에 반영합니다."""
    query = str(st.session_state.get("place_name", "") or "").strip()
    api_id, api_secret = get_naver_local_api_credentials()
    result = search_naver_local_api(query, api_id, api_secret, 5)
    items = result.get("items", []) or []
    st.session_state.naver_local_results = items
    if items:
        first_item = items[0]
        title = strip_html_tags(first_item.get("title", "")) or query
        road = strip_html_tags(first_item.get("roadAddress", ""))
        addr = strip_html_tags(first_item.get("address", ""))
        chosen_address = road or addr
        queue_location_update(title, chosen_address, "주소를 자동 입력했습니다.")
        st.session_state.naver_local_search_message = ""
    else:
        st.session_state.last_address_fetch_message = result.get("message", "장소 후보를 찾지 못했습니다.")
        st.session_state.naver_local_search_message = result.get("message", "")


def naver_item_map_link(item: dict, fallback_query: str) -> str:
    """선택 후보를 네이버 지도에서 열기 위한 링크를 만듭니다."""
    title = item.get("title") or fallback_query
    query = quote(str(title or "").strip())
    return f"https://map.naver.com/p/search/{query}?c=16.00,0,0,0,dh"

def capture_naver_map_region(place_query: str, wait_seconds: float = 8.0) -> tuple[bool, str]:
    """
    Playwright headless Chromium으로 네이버 지도 검색 URL을 백그라운드에서 열고,
    1920 x 1080 viewport 기준 좌표 (585, 87)~(1785, 1047)를 잘라 저장합니다.

    v23 변경점:
    - 큰 base64 문자열을 session_state에 직접 저장하지 않고 파일 경로를 저장합니다.
    - Playwright clip 캡처 대신 전체 화면 캡처 후 Pillow로 crop합니다.
      네이버 지도처럼 로딩이 늦거나 canvas/WebGL 렌더링이 있는 화면에서 더 안정적입니다.
    - 캡처 결과 파일을 좌측 입력부에서 바로 썸네일로 확인할 수 있도록 합니다.
    """
    query = str(place_query or "").strip()
    if not query:
        return False, "교육 장소명을 먼저 입력해 주세요."

    try:
        from playwright.sync_api import sync_playwright
    except Exception:
        return (
            False,
            "지도 자동 캡처에는 Playwright가 필요합니다. 터미널에서 `python -m pip install playwright` 실행 후, "
            "처음 한 번만 `python -m playwright install chromium`을 실행해 주세요.",
        )

    try:
        from PIL import Image, ImageStat
    except Exception:
        return (
            False,
            "지도 이미지를 좌표대로 자르려면 Pillow가 필요합니다. 터미널에서 `python -m pip install pillow`를 실행해 주세요.",
        )

    map_url = f"https://map.naver.com/p/search/{quote(query)}?c=15.00,0,0,0,dh"
    crop_box = (685, 167, 1685, 967)

    try:
        with sync_playwright() as pw:
            browser = pw.chromium.launch(
                headless=True,
                args=[
                    "--disable-blink-features=AutomationControlled",
                    "--no-sandbox",
                    "--disable-dev-shm-usage",
                    "--window-size=1920,1080",
                    "--force-device-scale-factor=1",
                    "--hide-scrollbars",
                    "--disable-gpu",
                ],
            )
            context = browser.new_context(
                viewport={"width": 1920, "height": 1080},
                screen={"width": 1920, "height": 1080},
                device_scale_factor=1,
                locale="ko-KR",
                timezone_id="Asia/Seoul",
                ignore_https_errors=True,
                user_agent=(
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/124.0.0.0 Safari/537.36"
                ),
            )
            page = context.new_page()
            page.goto(map_url, wait_until="load", timeout=90000)

            # 네이버 지도는 지도 타일/검색 결과가 늦게 붙는 경우가 많아서 조금 넉넉히 기다립니다.
            try:
                page.wait_for_load_state("networkidle", timeout=20000)
            except Exception:
                pass
            try:
                page.wait_for_selector("canvas, iframe, img, [class*=map], [class*=Map]", timeout=20000)
            except Exception:
                pass
            page.wait_for_timeout(int(max(4.0, float(wait_seconds)) * 1000))

            # 전체 viewport를 먼저 캡처한 뒤 지정 좌표로 crop합니다.
            screenshot_bytes = page.screenshot(type="png", full_page=False, animations="disabled")
            browser.close()

        image = Image.open(io.BytesIO(screenshot_bytes)).convert("RGB")
        width, height = image.size
        left, top, right, bottom = crop_box
        if width < right or height < bottom:
            return False, f"캡처 화면 크기가 예상보다 작습니다. 현재 캡처 크기: {width}x{height}, 필요 크기: 1785x1047"

        cropped = image.crop(crop_box)

        # 완전 백지/단색에 가까운 캡처인지 간단 점검합니다. 저장은 하되 안내 메시지에 표시합니다.
        stat = ImageStat.Stat(cropped.resize((120, 96)))
        avg_stddev = sum(stat.stddev) / len(stat.stddev)
        quality_note = ""
        if avg_stddev < 3:
            quality_note = " 다만 캡처 이미지가 거의 단색으로 감지됩니다. 네이버 지도 로딩/접근 차단 가능성이 있어 직접 첨부를 확인해 주세요."

        captures_dir = get_captures_dir()
        captures_dir.mkdir(parents=True, exist_ok=True)
        file_name = f"naver_map_capture_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"
        file_path = captures_dir / file_name
        cropped.save(file_path, format="PNG")

        st.session_state.captured_map_file_path = str(file_path)
        st.session_state.captured_map_data_url = ""
        st.session_state.last_captured_map_file = str(file_path)
        return True, f"지도 이미지를 자동 캡처해 적용했습니다.{quality_note}"
    except Exception as exc:
        return False, (
            "지도 자동 캡처에 실패했습니다. `python -m pip install playwright pillow` 및 "
            "`python -m playwright install chromium` 실행 여부를 확인해 주세요. "
            f"상세: {exc}"
        )



def _find_text_in_all_frames(page, selector: str, timeout_ms: int = 2500) -> str:
    """현재 페이지와 모든 iframe에서 selector의 첫 번째 텍스트를 찾습니다."""
    for frame in page.frames:
        try:
            locator = frame.locator(selector).first
            if locator.count() > 0:
                text = locator.inner_text(timeout=timeout_ms).strip()
                if text:
                    return text
        except Exception:
            continue
    return ""


def _wait_for_frame_by_name(page, frame_name: str, timeout_ms: int = 18000):
    """네이버 지도 iframe은 늦게 붙는 경우가 많아서 이름 기준으로 반복 탐색합니다."""
    deadline = time.time() + (timeout_ms / 1000)
    while time.time() < deadline:
        try:
            frame = page.frame(name=frame_name)
            if frame:
                return frame
        except Exception:
            pass
        try:
            page.wait_for_timeout(250)
        except Exception:
            time.sleep(0.25)
    return None


def _extract_korean_road_address(text: str) -> str:
    """네이버 지도 상세/검색 결과 텍스트에서 도로명 주소처럼 보이는 한 줄을 추출합니다."""
    raw = str(text or "")
    lines = [line.strip() for line in raw.splitlines() if line.strip()]
    province_pattern = r"(서울|부산|대구|인천|광주|대전|울산|세종|경기|강원|충북|충남|전북|전남|경북|경남|제주)"
    road_pattern = re.compile(province_pattern + r"[^\n]{0,90}?(?:로|길)\s*\d+(?:[\-\d]*)?(?:\s*[^\n]{0,45})?")
    for line in lines:
        if any(token in line for token in ["복사", "전화", "영업", "리뷰", "사진", "거리뷰", "출발", "도착"]):
            continue
        match = road_pattern.search(line)
        if match:
            address = match.group(0).strip()
            address = re.sub(r"\s+", " ", address)
            return address
    match = road_pattern.search(raw)
    if match:
        return re.sub(r"\s+", " ", match.group(0).strip())
    return ""


def _try_click_locator(locator, timeout_ms: int = 7000) -> bool:
    """React 이벤트 기반 href='#' 버튼 클릭을 여러 방식으로 시도합니다."""
    try:
        locator.wait_for(state="visible", timeout=timeout_ms)
    except Exception:
        return False

    try:
        locator.scroll_into_view_if_needed(timeout=timeout_ms)
    except Exception:
        pass

    click_methods = [
        lambda: locator.click(timeout=timeout_ms),
        lambda: locator.click(timeout=timeout_ms, force=True),
        lambda: locator.dispatch_event("click", timeout=timeout_ms),
        lambda: locator.evaluate("el => el.click()", timeout=timeout_ms),
        lambda: locator.press("Enter", timeout=timeout_ms),
    ]
    for method in click_methods:
        try:
            method()
            return True
        except Exception:
            continue
    return False


def _click_first_naver_place_result(page) -> tuple[bool, str]:
    """
    네이버 지도 검색 결과 첫 번째 항목을 클릭합니다.
    검색 결과의 a 태그 href는 대부분 '#': 실제 상세 이동은 React click handler로 처리됩니다.
    그래서 href를 읽지 않고, 제목/썸네일/role=button 요소를 실제 클릭합니다.
    """
    search_frame = _wait_for_frame_by_name(page, "searchIframe", timeout_ms=22000)
    candidate_frames = []
    if search_frame:
        candidate_frames.append(search_frame)
    candidate_frames.extend([frame for frame in page.frames if frame not in candidate_frames])

    selectors = [
        "#_pcmap_list_scroll_container > ul > li:nth-child(1) a.U70Fj",
        "#_pcmap_list_scroll_container > ul > li:nth-child(1) a.place_thumb",
        "#_pcmap_list_scroll_container > ul > li:nth-child(1) a[role='button']",
        "#_pcmap_list_scroll_container > ul > li:nth-child(1) [role='button']",
        "li.VLTHu:nth-child(1) a.U70Fj",
        "li.VLTHu:nth-child(1) a.place_thumb",
        "li:nth-child(1) a.U70Fj",
        "li:nth-child(1) a.place_thumb",
        "li:nth-child(1) [role='button']",
        "a.U70Fj:has(span.YwYLL)",
        "span.YwYLL",
    ]

    last_error = ""
    for frame in candidate_frames:
        for selector in selectors:
            try:
                locator = frame.locator(selector).first
                if locator.count() <= 0:
                    continue
                if _try_click_locator(locator):
                    return True, f"첫 번째 검색 결과를 클릭했습니다. selector={selector}"
            except Exception as exc:
                last_error = str(exc)
                continue

    # 최후 fallback: DOM 좌표로 첫 번째 결과 제목 위치를 클릭합니다.
    # searchIframe 안의 locator bounding_box는 메인 페이지 좌표로 변환되어 반환됩니다.
    try:
        if search_frame:
            locator = search_frame.locator("#_pcmap_list_scroll_container > ul > li:nth-child(1)").first
            locator.wait_for(state="visible", timeout=5000)
            box = locator.bounding_box(timeout=5000)
            if box:
                page.mouse.click(box["x"] + min(160, box["width"] / 2), box["y"] + min(42, box["height"] / 2))
                return True, "첫 번째 검색 결과를 좌표 클릭했습니다."
    except Exception as exc:
        last_error = str(exc)

    return False, f"첫 번째 검색 결과 클릭에 실패했습니다. {last_error}".strip()


def _extract_road_address_from_loaded_naver(page) -> str:
    """entryIframe 우선, 이후 전체 iframe에서 주소를 찾습니다."""
    target_frames = []
    entry_frame = _wait_for_frame_by_name(page, "entryIframe", timeout_ms=18000)
    if entry_frame:
        target_frames.append(entry_frame)
    target_frames.extend([frame for frame in page.frames if frame not in target_frames])

    selectors = [
        "#app-root div.place_section_content div.O8qbU.tQY7D div a span.pz7wy",
        "#app-root span.pz7wy",
        "span.pz7wy",
        "a span.pz7wy",
        "[class*='O8qbU'] [class*='pz7wy']",
    ]

    for _ in range(10):
        for frame in target_frames:
            for selector in selectors:
                try:
                    locator = frame.locator(selector).first
                    if locator.count() <= 0:
                        continue
                    text = locator.inner_text(timeout=1500).strip()
                    address = _extract_korean_road_address(text) or text
                    if address and ("로" in address or "길" in address):
                        return re.sub(r"\s+", " ", address.strip())
                except Exception:
                    continue
        try:
            page.wait_for_timeout(700)
        except Exception:
            time.sleep(0.7)

    for frame in target_frames:
        try:
            body_text = frame.locator("body").inner_text(timeout=2500)
            address = _extract_korean_road_address(body_text)
            if address:
                return address
        except Exception:
            continue

    return ""


def fetch_naver_road_address(place_query: str, wait_seconds: float = 5.0) -> tuple[bool, str]:
    """
    네이버 지도 검색 결과의 첫 번째 장소 상세 페이지에서 도로명 주소를 가져옵니다.
    v25 변경점:
    - href='#'를 링크로 해석하지 않고 실제 React click handler를 실행합니다.
    - 썸네일보다 제목 링크(a.U70Fj)를 우선 클릭합니다.
    - searchIframe/entryIframe 로딩을 더 오래 기다립니다.
    - 실패 시 디버그용 스크린샷/HTML을 assets/captures에 저장합니다.
    """
    query = str(place_query or "").strip()
    if not query:
        return False, "교육 장소명을 먼저 입력해 주세요."

    try:
        from playwright.sync_api import sync_playwright
    except Exception:
        return False, "도로명 주소 가져오기에는 Playwright가 필요합니다. `python -m pip install playwright` 및 `python -m playwright install chromium`을 실행해 주세요."

    map_url = f"https://map.naver.com/p/search/{quote(query)}?c=16.00,0,0,0,dh"
    debug_dir = get_captures_dir()
    debug_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    try:
        with sync_playwright() as pw:
            browser = pw.chromium.launch(
                headless=True,
                args=[
                    "--disable-blink-features=AutomationControlled",
                    "--no-sandbox",
                    "--disable-dev-shm-usage",
                    "--window-size=1920,1080",
                    "--force-device-scale-factor=1",
                    "--hide-scrollbars",
                    "--disable-gpu",
                ],
            )
            context = browser.new_context(
                viewport={"width": 1920, "height": 1080},
                screen={"width": 1920, "height": 1080},
                device_scale_factor=1,
                locale="ko-KR",
                timezone_id="Asia/Seoul",
                ignore_https_errors=True,
                user_agent=(
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/124.0.0.0 Safari/537.36"
                ),
            )
            page = context.new_page()
            page.goto(map_url, wait_until="domcontentloaded", timeout=70000)
            page.wait_for_timeout(int(max(3.0, float(wait_seconds)) * 1000))

            clicked, click_message = _click_first_naver_place_result(page)
            if not clicked:
                try:
                    page.screenshot(path=str(debug_dir / f"address_click_failed_{stamp}.png"), full_page=True)
                except Exception:
                    pass
                try:
                    (debug_dir / f"address_click_failed_{stamp}.html").write_text(page.content(), encoding="utf-8")
                except Exception:
                    pass
                browser.close()
                return False, click_message

            # entryIframe이 붙고 상세 DOM이 그려질 시간을 줍니다.
            page.wait_for_timeout(2500)
            address = _extract_road_address_from_loaded_naver(page)

            if not address:
                try:
                    page.screenshot(path=str(debug_dir / f"address_extract_failed_{stamp}.png"), full_page=True)
                except Exception:
                    pass
                try:
                    frames_dump = []
                    for frame in page.frames:
                        frame_name = getattr(frame, "name", "") or "no_name"
                        frame_url = getattr(frame, "url", "") or ""
                        try:
                            frame_text = frame.locator("body").inner_text(timeout=1000)[:4000]
                        except Exception:
                            frame_text = ""
                        frames_dump.append(f"\n--- FRAME: {frame_name} ---\nURL: {frame_url}\n{frame_text}")
                    (debug_dir / f"address_extract_failed_{stamp}.txt").write_text("\n".join(frames_dump), encoding="utf-8")
                except Exception:
                    pass

            browser.close()

        if address:
            return True, address.strip()
        return False, (
            "첫 번째 검색 결과 상세 화면까지는 클릭했지만 도로명 주소를 찾지 못했습니다. "
            "assets/captures 폴더의 address_extract_failed 파일을 확인해 주세요."
        )
    except Exception as exc:
        return False, f"도로명 주소 가져오기에 실패했습니다. 상세: {exc}"


def fetch_road_address_callback() -> None:
    ok, result = fetch_naver_road_address(st.session_state.get("place_name", ""))
    if ok:
        st.session_state.road_address = result
        # 가져온 주소는 검색용 주소 필드와 안내문 표시 주소 필드에 같이 넣어줍니다.
        # 이후 사용자가 안내문 표시 주소만 별도로 수정할 수 있습니다.
        st.session_state.display_road_address = result
        place_for_display = str(st.session_state.get("place_name", "") or "").strip()
        if not str(st.session_state.get("display_place_name", "")).strip():
            st.session_state.display_place_name = place_for_display
        if place_for_display:
            st.session_state.display_location_text = f"{place_for_display}  ({result})"
        else:
            st.session_state.display_location_text = result
        st.session_state.last_address_fetch_message = f"도로명 주소를 가져왔습니다: {result}"
    else:
        st.session_state.last_address_fetch_message = result



def capture_naver_map_variants(place_query: str, wait_seconds: float = 0.8) -> dict:
    """
    빠른 지도 캡처 전용 경로입니다.
    - 주소 추출은 하지 않습니다.
    - 네이버 지도 페이지를 한 번만 열고 viewport 한 장만 캡처합니다.
    - 두 크기 이미지는 같은 원본 스크린샷에서 crop 합니다.
    - 대기 시간은 짧게 두고, 실패하면 빠르게 종료합니다.
    """
    query = str(place_query or "").strip()
    if not query:
        return {"ok": False, "message": "교육 장소명 / 검색어를 먼저 입력해 주세요.", "map_file_paths": {}}

    try:
        from playwright.sync_api import sync_playwright
    except Exception:
        return {"ok": False, "message": "지도 캡처에는 Playwright가 필요합니다. `python -m pip install playwright` 및 `python -m playwright install chromium`을 실행해 주세요.", "map_file_paths": {}}

    try:
        from PIL import Image, ImageStat
    except Exception:
        return {"ok": False, "message": "지도 이미지를 자르려면 Pillow가 필요합니다. `python -m pip install pillow`를 실행해 주세요.", "map_file_paths": {}}

    map_url = f"https://map.naver.com/p/search/{quote(query)}?c=16.00,0,0,0,dh"
    crop_boxes = get_capture_boxes()
    captures_dir = get_captures_dir()
    captures_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    try:
        with sync_playwright() as pw:
            browser = pw.chromium.launch(
                headless=True,
                args=[
                    "--disable-blink-features=AutomationControlled",
                    "--no-sandbox",
                    "--disable-dev-shm-usage",
                    "--window-size=1920,1080",
                    "--force-device-scale-factor=1",
                    "--hide-scrollbars",
                    "--disable-gpu",
                ],
            )
            context = browser.new_context(
                viewport={"width": 1920, "height": 1080},
                screen={"width": 1920, "height": 1080},
                device_scale_factor=1,
                locale="ko-KR",
                timezone_id="Asia/Seoul",
                ignore_https_errors=True,
                user_agent=(
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/124.0.0.0 Safari/537.36"
                ),
            )
            page = context.new_page()
            page.goto(map_url, wait_until="domcontentloaded", timeout=15000)
            try:
                page.wait_for_selector("iframe#searchIframe, canvas, [class*=map], [class*=Map], img", timeout=3500)
            except Exception:
                pass
            page.wait_for_timeout(int(max(0.4, float(wait_seconds)) * 1000))
            screenshot_bytes = page.screenshot(type="png", full_page=False, animations="disabled")
            browser.close()

        image = Image.open(io.BytesIO(screenshot_bytes)).convert("RGB")
        width, height = image.size
        map_file_paths = {}
        blank_notes = []

        for variant_name, crop_box in crop_boxes.items():
            left, top, right, bottom = crop_box
            if width < right or height < bottom:
                continue
            cropped = image.crop(crop_box)
            stat = ImageStat.Stat(cropped.resize((120, 96)))
            avg_stddev = sum(stat.stddev) / len(stat.stddev)
            safe_name = "small" if "작은" in variant_name else "normal"
            file_path = captures_dir / f"naver_map_capture_{safe_name}_{stamp}.png"
            cropped.save(file_path, format="PNG")
            map_file_paths[variant_name] = str(file_path)
            if avg_stddev < 2.5:
                blank_notes.append(variant_name)

        if not map_file_paths:
            return {"ok": False, "message": f"지도 캡처 화면 크기가 예상보다 작습니다. 현재 {width}x{height}입니다.", "map_file_paths": {}}

        message = "지도 이미지를 가져왔습니다."
        if blank_notes:
            message += " 일부 캡처가 단색에 가깝습니다. 필요하면 직접 첨부 이미지를 사용해 주세요."
        return {"ok": True, "message": message, "map_file_paths": map_file_paths}
    except Exception as exc:
        return {"ok": False, "message": f"지도 이미지 가져오기에 실패했습니다. 상세: {exc}", "map_file_paths": {}}


def fetch_naver_road_address_result(place_query: str, wait_seconds: float = 2.5) -> dict:
    """주소만 별도로 가져옵니다. 실패해도 지도 캡처에는 영향을 주지 않도록 dict로 반환합니다."""
    ok, result = fetch_naver_road_address(place_query, wait_seconds=wait_seconds)
    if ok:
        query = str(place_query or "").strip()
        return {
            "ok": True,
            "road_address": result.strip(),
            "display_location_text": format_location_line(query, result),
            "message": f"도로명 주소를 가져왔습니다: {result.strip()}",
        }
    return {"ok": False, "road_address": "", "display_location_text": "", "message": result}


def start_naver_map_background_job() -> None:
    query = str(st.session_state.get("place_name", "") or "").strip()
    if not query:
        st.session_state.naver_map_background_message = "교육 장소명 / 검색어를 먼저 입력해 주세요."
        return
    future = st.session_state.get("naver_map_future")
    if future is not None and not future.done():
        st.session_state.naver_map_background_message = "이미 지도 이미지 가져오기가 실행 중입니다. 잠시 후 결과를 확인해 주세요."
        return
    executor = get_naver_executor()
    st.session_state.naver_map_future = executor.submit(capture_naver_map_variants, query, 0.8)
    st.session_state.naver_map_started_at = time.time()
    st.session_state.naver_map_background_message = "지도 이미지를 가져오는 중입니다."


def start_naver_address_background_job() -> None:
    query = str(st.session_state.get("place_name", "") or "").strip()
    if not query:
        st.session_state.naver_address_background_message = "교육 장소명 / 검색어를 먼저 입력해 주세요."
        return
    future = st.session_state.get("naver_address_future")
    if future is not None and not future.done():
        st.session_state.naver_address_background_message = "이미 도로명 주소 가져오기가 실행 중입니다. 잠시 후 결과를 확인해 주세요."
        return
    executor = get_naver_executor()
    st.session_state.naver_address_future = executor.submit(fetch_naver_road_address_result, query, 2.5)
    st.session_state.naver_address_started_at = time.time()
    st.session_state.naver_address_background_message = "도로명 주소 가져오기를 백그라운드에서 시작했습니다. 다른 입력 작업을 계속할 수 있습니다."


def poll_split_naver_jobs() -> None:
    map_future = st.session_state.get("naver_map_future")
    if map_future is not None:
        if not map_future.done():
            started = st.session_state.get("naver_map_started_at", time.time())
            elapsed = int(time.time() - started)
            st.session_state.naver_map_background_message = f"지도 이미지 가져오는 중... {elapsed}초"
        else:
            try:
                result = map_future.result()
            except Exception as exc:
                result = {"ok": False, "message": f"지도 이미지 백그라운드 작업 오류: {exc}", "map_file_paths": {}}
            st.session_state.pop("naver_map_future", None)
            if result.get("map_file_paths"):
                st.session_state.captured_map_files = result["map_file_paths"]
                selected = st.session_state.get("selected_map_capture_variant", "기본 크기 (1000×800 / 5:4)")
                chosen_path = result["map_file_paths"].get(selected) or result["map_file_paths"].get("기본 크기 (1000×800 / 5:4)") or next(iter(result["map_file_paths"].values()))
                st.session_state.captured_map_file_path = chosen_path
                st.session_state.captured_map_data_url = ""
                st.session_state.last_captured_map_file = chosen_path
                st.session_state.last_map_capture_message = "지도 이미지를 기본 크기(1000×800, 5:4)로 가져왔습니다."
            st.session_state.naver_map_background_message = result.get("message", "지도 이미지 작업이 완료되었습니다.")

    address_future = st.session_state.get("naver_address_future")
    if address_future is not None:
        if not address_future.done():
            started = st.session_state.get("naver_address_started_at", time.time())
            elapsed = int(time.time() - started)
            st.session_state.naver_address_background_message = f"도로명 주소 가져오는 중... 약 {elapsed}초 경과."
        else:
            try:
                result = address_future.result()
            except Exception as exc:
                result = {"ok": False, "message": f"도로명 주소 백그라운드 작업 오류: {exc}", "road_address": "", "display_location_text": ""}
            st.session_state.pop("naver_address_future", None)
            if result.get("road_address"):
                st.session_state.road_address = result["road_address"]
                st.session_state.last_address_fetch_message = f"도로명 주소를 가져왔습니다: {result['road_address']}"
            if result.get("display_location_text"):
                st.session_state.display_location_text = result["display_location_text"]
            st.session_state.naver_address_background_message = result.get("message", "도로명 주소 작업이 완료되었습니다.")


@st.cache_resource
def get_naver_executor() -> ThreadPoolExecutor:
    """네이버 지도/주소 자동화 작업을 백그라운드에서 실행합니다."""
    return ThreadPoolExecutor(max_workers=2)


def get_capture_boxes() -> dict[str, tuple[int, int, int, int]]:
    # 1000 x 800 (5:4)
    # 중앙: (1185, 567)
    # 좌측상단: (685, 167) / 우측상단: (1685, 167)
    # 좌측하단: (685, 967) / 우측하단: (1685, 967)
    return {
        "기본 크기 (1000×800 / 5:4)": (685, 167, 1685, 967),
    }


def format_location_line(place: str, address: str) -> str:
    clean_place = str(place or "").strip()
    clean_address = str(address or "").strip()
    if clean_place and clean_address:
        return f"{clean_place}  ({clean_address})"
    return clean_place or clean_address


def capture_naver_map_and_address(place_query: str, wait_seconds: float = 3.0) -> dict:
    """
    네이버 지도 검색 URL을 한 번만 열어서 지도 이미지 2종(작은/일반)을 먼저 캡처하고,
    그 다음 첫 번째 결과를 클릭해 도로명 주소를 가져옵니다.
    """
    query = str(place_query or "").strip()
    if not query:
        return {"ok": False, "message": "교육 장소명을 먼저 입력해 주세요."}

    try:
        from playwright.sync_api import sync_playwright
    except Exception:
        return {"ok": False, "message": "Playwright가 필요합니다. `python -m pip install playwright` 및 `python -m playwright install chromium`을 실행해 주세요."}

    try:
        from PIL import Image, ImageStat
    except Exception:
        return {"ok": False, "message": "Pillow가 필요합니다. `python -m pip install pillow`를 실행해 주세요."}

    map_url = f"https://map.naver.com/p/search/{quote(query)}?c=16.00,0,0,0,dh"
    crop_boxes = get_capture_boxes()
    debug_dir = get_captures_dir()
    debug_dir.mkdir(parents=True, exist_ok=True)
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    result = {
        "ok": False,
        "message": "",
        "map_file_path": "",
        "map_file_paths": {},
        "road_address": "",
        "display_location_text": "",
    }

    try:
        with sync_playwright() as pw:
            browser = pw.chromium.launch(
                headless=True,
                args=[
                    "--disable-blink-features=AutomationControlled",
                    "--no-sandbox",
                    "--disable-dev-shm-usage",
                    "--window-size=1920,1080",
                    "--force-device-scale-factor=1",
                    "--hide-scrollbars",
                    "--disable-gpu",
                ],
            )
            context = browser.new_context(
                viewport={"width": 1920, "height": 1080},
                screen={"width": 1920, "height": 1080},
                device_scale_factor=1,
                locale="ko-KR",
                timezone_id="Asia/Seoul",
                ignore_https_errors=True,
                user_agent=(
                    "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                    "AppleWebKit/537.36 (KHTML, like Gecko) "
                    "Chrome/124.0.0.0 Safari/537.36"
                ),
            )
            page = context.new_page()
            page.goto(map_url, wait_until="domcontentloaded", timeout=50000)

            # 지도/검색 UI가 최소한 붙을 때까지만 기다립니다. 오래 기다리는 networkidle은 피합니다.
            try:
                page.wait_for_selector("iframe#searchIframe, canvas, [class*=map], [class*=Map]", timeout=12000)
            except Exception:
                pass
            page.wait_for_timeout(int(max(2.0, float(wait_seconds)) * 1000))

            # 1) 상세 클릭 전, 지도 화면을 먼저 1회 캡처한 뒤 두 크기로 잘라 저장합니다.
            screenshot_bytes = page.screenshot(type="png", full_page=False, animations="disabled")
            image = Image.open(io.BytesIO(screenshot_bytes)).convert("RGB")
            width, height = image.size
            blank_notes = []
            for variant_name, crop_box in crop_boxes.items():
                left, top, right, bottom = crop_box
                if width >= right and height >= bottom:
                    cropped = image.crop(crop_box)
                    stat = ImageStat.Stat(cropped.resize((120, 96)))
                    avg_stddev = sum(stat.stddev) / len(stat.stddev)
                    safe_name = "small" if "작은" in variant_name else "normal"
                    file_name = f"naver_map_capture_{safe_name}_{stamp}.png"
                    file_path = debug_dir / file_name
                    cropped.save(file_path, format="PNG")
                    result["map_file_paths"][variant_name] = str(file_path)
                    if not result["map_file_path"] and "일반" in variant_name:
                        result["map_file_path"] = str(file_path)
                    if avg_stddev < 3:
                        blank_notes.append(variant_name)
                else:
                    result["message"] += f"{variant_name} 캡처 화면 크기가 부족합니다. 현재 {width}x{height}. "
            if not result["map_file_path"] and result["map_file_paths"]:
                result["map_file_path"] = next(iter(result["map_file_paths"].values()))
            if blank_notes:
                result["message"] += "일부 지도 이미지가 거의 단색으로 감지됩니다: " + ", ".join(blank_notes) + ". "

            # 2) 지도 캡처 이후 첫 번째 검색 결과 클릭 → 상세 주소 추출
            clicked, click_message = _click_first_naver_place_result(page)
            if clicked:
                page.wait_for_timeout(1600)
                address = _extract_road_address_from_loaded_naver(page)
                if address:
                    result["road_address"] = address.strip()
                    result["display_location_text"] = format_location_line(query, address)
                else:
                    try:
                        frames_dump = []
                        for frame in page.frames:
                            frame_name = getattr(frame, "name", "") or "no_name"
                            frame_url = getattr(frame, "url", "") or ""
                            try:
                                frame_text = frame.locator("body").inner_text(timeout=800)[:2500]
                            except Exception:
                                frame_text = ""
                            frames_dump.append(f"\n--- FRAME: {frame_name} ---\nURL: {frame_url}\n{frame_text}")
                        (debug_dir / f"combined_address_extract_failed_{stamp}.txt").write_text("\n".join(frames_dump), encoding="utf-8")
                    except Exception:
                        pass
                    result["message"] += "상세 화면에서 도로명 주소를 찾지 못했습니다. "
            else:
                result["message"] += click_message + " "

            browser.close()

        if result["map_file_paths"] or result["road_address"]:
            result["ok"] = True
            details = []
            if result["map_file_paths"]:
                details.append("지도 2종 캡처 완료")
            if result["road_address"]:
                details.append(f"주소 가져오기 완료: {result['road_address']}")
            result["message"] = (" / ".join(details) + (" " + result["message"] if result["message"] else "")).strip()
        else:
            result["ok"] = False
            result["message"] = result["message"] or "지도 캡처와 도로명 주소 가져오기에 모두 실패했습니다."
        return result
    except Exception as exc:
        return {"ok": False, "message": f"지도/주소 동시 가져오기에 실패했습니다. 상세: {exc}"}


def start_naver_background_job() -> None:
    query = str(st.session_state.get("place_name", "") or "").strip()
    if not query:
        st.session_state.naver_background_message = "교육 장소명 / 검색어를 먼저 입력해 주세요."
        return
    future = st.session_state.get("naver_combined_future")
    if future is not None and not future.done():
        st.session_state.naver_background_message = "이미 지도/주소 가져오기가 실행 중입니다. 잠시 후 결과를 확인해 주세요."
        return
    executor = get_naver_executor()
    st.session_state.naver_combined_future = executor.submit(
        capture_naver_map_and_address,
        query,
        2.6,
    )
    st.session_state.naver_background_message = "지도 캡처와 도로명 주소 가져오기를 백그라운드에서 시작했습니다. 다른 입력값을 계속 수정해도 됩니다."
    st.session_state.naver_combined_started_at = time.time()


def poll_naver_background_job() -> None:
    future = st.session_state.get("naver_combined_future")
    if future is None:
        return
    if not future.done():
        started = st.session_state.get("naver_combined_started_at", time.time())
        elapsed = int(time.time() - started)
        st.session_state.naver_background_message = f"지도/주소 가져오는 중... 약 {elapsed}초 경과. 다른 입력 작업을 계속할 수 있습니다."
        return
    try:
        result = future.result()
    except Exception as exc:
        result = {"ok": False, "message": f"백그라운드 작업 오류: {exc}"}
    st.session_state.pop("naver_combined_future", None)
    if result.get("map_file_paths"):
        st.session_state.captured_map_files = result["map_file_paths"]
        selected = st.session_state.get("selected_map_capture_variant", "기본 크기 (1000×800 / 5:4)")
        chosen_path = result["map_file_paths"].get(selected) or result.get("map_file_path") or next(iter(result["map_file_paths"].values()))
        st.session_state.captured_map_file_path = chosen_path
        st.session_state.captured_map_data_url = ""
        st.session_state.last_captured_map_file = chosen_path
        st.session_state.last_map_capture_message = "지도 이미지를 기본 크기(1000×800, 5:4)로 자동 캡처했습니다."
    elif result.get("map_file_path"):
        st.session_state.captured_map_file_path = result["map_file_path"]
        st.session_state.captured_map_data_url = ""
        st.session_state.last_captured_map_file = result["map_file_path"]
        st.session_state.last_map_capture_message = "지도 이미지를 자동 캡처해 미리보기에 적용했습니다."
    if result.get("road_address"):
        st.session_state.road_address = result["road_address"]
        st.session_state.last_address_fetch_message = f"도로명 주소를 가져왔습니다: {result['road_address']}"
    if result.get("display_location_text"):
        st.session_state.display_location_text = result["display_location_text"]
    st.session_state.naver_background_message = result.get("message", "작업이 완료되었습니다.")


def is_valid_hex_color(value: str) -> bool:
    return bool(re.fullmatch(r"#[0-9A-Fa-f]{6}", str(value or "").strip()))


def get_query_param(name: str) -> str:
    try:
        value = st.query_params.get(name, "")
        if isinstance(value, list):
            return str(value[0]) if value else ""
        return str(value or "")
    except Exception:
        try:
            params = st.experimental_get_query_params()
            value = params.get(name, [""])
            return str(value[0]) if value else ""
        except Exception:
            return ""


def css_font_stack(primary_font: str) -> str:
    clean = str(primary_font or "").strip()
    clean = clean.replace("'", "").replace('"', "").replace(";", "").replace("\\", "")
    if not clean:
        clean = "Malgun Gothic"
    return f"'{clean}', 'Malgun Gothic', 'Apple SD Gothic Neo', Arial, sans-serif"


FONT_EXTENSIONS = (".ttf", ".otf", ".woff", ".woff2")
FONT_MIME_TYPES = {
    ".ttf": "font/ttf",
    ".otf": "font/otf",
    ".woff": "font/woff",
    ".woff2": "font/woff2",
}
FONT_FORMATS = {
    ".ttf": "truetype",
    ".otf": "opentype",
    ".woff": "woff",
    ".woff2": "woff2",
}


def get_fonts_dir() -> Path:
    return Path(__file__).resolve().parent / "assets" / "fonts"


@st.cache_data(show_spinner=False)
def list_font_files(fonts_dir_text: str) -> list[str]:
    fonts_dir = Path(fonts_dir_text)
    if not fonts_dir.exists() or not fonts_dir.is_dir():
        return []
    files: list[str] = []
    for ext in FONT_EXTENSIONS:
        files.extend(str(path) for path in fonts_dir.glob(f"*{ext}"))
        files.extend(str(path) for path in fonts_dir.glob(f"*{ext.upper()}"))
    return sorted(set(files), key=lambda x: Path(x).name.lower())


def font_family_from_file(font_path: str) -> str:
    name = Path(font_path).stem.strip()
    return re.sub(r"[^0-9A-Za-z가-힣 _.-]", "", name) or "CustomFont"


def sanitize_file_name(value: str) -> str:
    name = str(value or "").strip()
    name = re.sub(r'[\\/:*?"<>|]+', "_", name)
    name = re.sub(r"\s+", "_", name)
    name = name.strip("._ ")
    return name or "education_notice_photo_style"


def _safe_widget_key_part(value: object) -> str:
    """동적 입력 위젯 key에 쓸 수 있도록 컬럼명을 정리합니다."""
    text_value = str(value or "").strip()
    text_value = re.sub(r"[^0-9A-Za-z가-힣_]+", "_", text_value)
    text_value = text_value.strip("_")
    return text_value or "column"


def _normalize_editor_rows(records: object, columns: list[str], min_rows: int = 1) -> list[dict]:
    normalized_rows: list[dict] = []
    for row in to_records(records):
        normalized_row = {column: _cell_value_by_column(row, column) for column in columns}
        if any(str(value or "").strip() for value in normalized_row.values()):
            normalized_rows.append(normalized_row)
    while len(normalized_rows) < max(1, int(min_rows or 1)):
        normalized_rows.append({column: "" for column in columns})
    return normalized_rows


def render_white_theme_grid_editor(
    records: object,
    columns: list[str],
    key_prefix: str,
    column_labels: dict[str, str] | None = None,
    min_rows: int = 1,
    add_button_label: str = "행 추가",
) -> list[dict]:
    """화이트 모드 전용 단순 편집기입니다.

    st.data_editor는 내부 셀을 canvas로 렌더링해서, 앱 기본 테마가 다크로 고정된 경우
    CSS만으로 셀 배경/글자색을 안정적으로 바꾸기 어렵습니다. 화이트 모드에서는
    흰 배경 + 검정 글씨의 기본 입력창 기반 편집기를 사용합니다.
    """
    column_labels = column_labels or {}
    rows_key = f"white_grid_{key_prefix}_rows"

    if rows_key not in st.session_state:
        st.session_state[rows_key] = _normalize_editor_rows(records, columns, min_rows=min_rows)

    rows = st.session_state.get(rows_key, [])
    if not isinstance(rows, list):
        rows = []
    rows = _normalize_editor_rows(rows, columns, min_rows=min_rows)

    st.markdown(
        """
        <style>
        .white-grid-editor {
            border: 1px solid #c9c9c9;
            border-radius: 12px;
            background: #FFFFFF;
            padding: 10px 10px 8px 10px;
            margin: 6px 0 10px 0;
        }
        .white-grid-header {
            color: #111827;
            font-size: 12px;
            line-height: 18px;
            font-weight: 850;
            padding: 0 2px 4px 2px;
        }
        .white-grid-editor div[data-baseweb="input"] > div,
        .white-grid-editor input {
            background: #FFFFFF !important;
            color: #111827 !important;
            border-color: #c9c9c9 !important;
            caret-color: #111827 !important;
        }
        .white-grid-editor input::placeholder {
            color: #9CA3AF !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    st.markdown('<div class="white-grid-editor">', unsafe_allow_html=True)

    if len(columns) == 1:
        header_cols = st.columns([1, 0.13])
    else:
        weights = [1.8 if any(token in column for token in ["내용", "주제", "과정", "비고", "메모", "role"]) else 1 for column in columns]
        header_cols = st.columns(weights + [0.13])
    for idx, column in enumerate(columns):
        header_cols[idx].markdown(
            f'<div class="white-grid-header">{esc(column_labels.get(column, column))}</div>',
            unsafe_allow_html=True,
        )
    header_cols[-1].markdown('<div class="white-grid-header">삭제</div>', unsafe_allow_html=True)

    edited_rows: list[dict] = []
    delete_index: int | None = None
    for row_idx, row in enumerate(rows):
        if len(columns) == 1:
            row_cols = st.columns([1, 0.13], vertical_alignment="center")
        else:
            weights = [1.8 if any(token in column for token in ["내용", "주제", "과정", "비고", "메모", "role"]) else 1 for column in columns]
            row_cols = st.columns(weights + [0.13], vertical_alignment="center")
        new_row: dict[str, str] = {}
        for col_idx, column in enumerate(columns):
            safe_column = _safe_widget_key_part(column)
            widget_key = f"white_grid_{key_prefix}_{row_idx}_{safe_column}"
            if widget_key not in st.session_state:
                st.session_state[widget_key] = str(row.get(column, "") or "")
            new_row[column] = row_cols[col_idx].text_input(
                column_labels.get(column, column),
                key=widget_key,
                label_visibility="collapsed",
                placeholder=column_labels.get(column, column),
            )
        if row_cols[-1].button("×", key=f"white_grid_{key_prefix}_{row_idx}_delete", help="이 행 삭제"):
            delete_index = row_idx
        edited_rows.append(new_row)

    st.markdown('</div>', unsafe_allow_html=True)

    if delete_index is not None:
        edited_rows.pop(delete_index)
        if not edited_rows:
            edited_rows = [{column: "" for column in columns}]
        st.session_state[rows_key] = edited_rows
        for key in list(st.session_state.keys()):
            if key.startswith(f"white_grid_{key_prefix}_") and key != rows_key:
                st.session_state.pop(key, None)
        st.rerun()

    action_cols = st.columns([1, 1])
    if action_cols[0].button(add_button_label, key=f"white_grid_{key_prefix}_add", use_container_width=True):
        edited_rows.append({column: "" for column in columns})
        st.session_state[rows_key] = edited_rows
        st.rerun()
    if action_cols[1].button("빈 행 정리", key=f"white_grid_{key_prefix}_clean", use_container_width=True):
        cleaned = [row for row in edited_rows if any(str(value or "").strip() for value in row.values())]
        if not cleaned:
            cleaned = [{column: "" for column in columns}]
        st.session_state[rows_key] = cleaned
        for key in list(st.session_state.keys()):
            if key.startswith(f"white_grid_{key_prefix}_") and key != rows_key:
                st.session_state.pop(key, None)
        st.rerun()

    st.session_state[rows_key] = edited_rows
    return edited_rows


def build_download_html_document(final_mail_html: str, font_stack: str, page_title: str = "교육 안내문") -> str:
    """브라우저에서 HTML 파일을 직접 열어도 한글이 깨지지 않도록 완성형 문서로 감쌉니다."""
    safe_title = html.escape(str(page_title or "교육 안내문"), quote=False)
    return f"""<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <meta http-equiv="Content-Type" content="text/html; charset=UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{safe_title}</title>
    <style>
        body {{
            margin: 0;
            padding: 32px 0;
            background: #EEF2F6;
            font-family: {font_stack};
        }}
    </style>
</head>
<body>
{final_mail_html}
</body>
</html>"""


def build_embedded_font_css(font_path: str, font_family: str) -> str:
    path = Path(font_path)
    if not path.exists() or path.suffix.lower() not in FONT_EXTENSIONS:
        return ""
    try:
        font_bytes = path.read_bytes()
    except OSError:
        return ""
    suffix = path.suffix.lower()
    mime = FONT_MIME_TYPES.get(suffix, "font/ttf")
    font_format = FONT_FORMATS.get(suffix, "truetype")
    data_url = f"data:{mime};base64,{base64.b64encode(font_bytes).decode('ascii')}"
    safe_family = font_family.replace("'", "").replace('"', "").replace(";", "")
    return f"""
    <style>
    @font-face {{
        font-family: '{safe_family}';
        src: url('{data_url}') format('{font_format}');
        font-weight: 400 900;
        font-style: normal;
        font-display: swap;
    }}
    </style>
    """


HIGHLIGHT_LABELS = {
    "title": "제목/교육명",
    "welcome": "상단 안내 문구",
    "time": "교육 시작시간 강조 문구",
    "brand": "브랜드 컬러",
    "header_text": "컬러박스 내 글자색",
    "logo": "로고",
    "font": "폰트",
    "overview": "교육 개요",
    "location": "교육 장소",
    "curriculum": "커리큘럼",
    "notice": "안내 사항",
    "contact": "관련 문의",
}


def set_active_highlight(zone: str) -> None:
    # 입력 중인 영역 하이라이트 기능은 v16에서 제거했습니다.
    return None


def render_template_html(template: str, replacements: dict[str, str]) -> str:
    """사용자가 수정한 문구의 플레이스홀더만 HTML로 치환하고 나머지는 안전하게 escape합니다."""
    template = str(template or "")
    tokens = sorted(replacements.keys(), key=len, reverse=True)
    result = ""
    idx = 0
    while idx < len(template):
        matched_token = None
        for token in tokens:
            if template.startswith(token, idx):
                matched_token = token
                break
        if matched_token is not None:
            result += replacements[matched_token]
            idx += len(matched_token)
            continue

        char = template[idx]
        if char == "\n":
            result += "<br>"
        else:
            result += esc(char)
        idx += 1
    return result


# -----------------------------
# Streamlit UI 스타일
# -----------------------------
init_defaults()
migrate_curriculum_day_defaults()

# 화면 테마는 상단 우측 셀렉트박스에서 바꿉니다.
# CSS는 위젯 렌더링 전 현재 session_state 값을 기준으로 먼저 적용합니다.
app_theme = st.session_state.get("app_theme", DEFAULT_VALUES["app_theme"])

st.markdown(
    """
    <style>
    :root {
        --app-bg: #050505;
        --card-bg: #1B1B1B;
        --card-bg-strong: #222222;
        --field-bg: #2B2B2B;
        --field-bg-focus: #333333;
        --line: #3A3A3A;
        --line-strong: #5A5A5A;
        --text: #F7F7F7;
        --muted: #B8B8B8;
        --soft: #242424;
        --button-bg: #FFFFFF;
        --button-text: #111111;
    }

    .stApp {
        background: var(--app-bg);
        color: var(--text);
    }

    .block-container {
        padding-top: 2.2rem;
        padding-bottom: 3rem;
        max-width: 1540px;
    }

    .main-title {
        font-size: 30px;
        font-weight: 850;
        color: var(--text);
        margin-bottom: 4px;
        letter-spacing: -0.7px;
    }

    .sub-title {
        font-size: 14px;
        color: var(--muted);
        margin-bottom: 22px;
    }

    .theme-select-note {
        color: var(--muted);
        font-size: 12px;
        line-height: 18px;
        text-align: right;
        margin-top: 28px;
        margin-bottom: 4px;
        font-weight: 800;
    }

    .st-key-app_theme {
        position: relative;
        z-index: 1;
        margin-bottom: 10px;
    }

    /* 카드형 입력/출력 섹션 */
    [data-testid="stVerticalBlockBorderWrapper"] {
        background: var(--card-bg);
        border: 1px solid var(--line);
        border-radius: 18px;
        box-shadow: 0 18px 42px rgba(0, 0, 0, 0.35);
        padding: 6px 6px;
    }

    [data-testid="stVerticalBlockBorderWrapper"] [data-testid="stMarkdownContainer"] h4,
    [data-testid="stVerticalBlockBorderWrapper"] [data-testid="stMarkdownContainer"] h3 {
        color: var(--text);
        letter-spacing: -0.35px;
        margin-bottom: 0.25rem;
    }

    /* 일반 텍스트/라벨 */
    [data-testid="stMarkdownContainer"],
    [data-testid="stMarkdownContainer"] p,
    [data-testid="stMarkdownContainer"] li,
    label,
    .stRadio label,
    .stCheckbox label,
    .stSlider label,
    .stColorPicker label,
    .stFileUploader label {
        color: var(--text) !important;
    }

    [data-testid="stCaptionContainer"],
    .mini-help,
    small,
    .st-emotion-cache-1wmy9hl,
    .st-emotion-cache-ue6h4q {
        color: var(--muted) !important;
    }

    /* 입력창: 차콜 배경 */
    div[data-baseweb="input"] > div,
    div[data-baseweb="textarea"] > div,
    div[data-baseweb="select"] > div {
        border-radius: 12px !important;
        border-color: var(--line-strong) !important;
        background-color: var(--field-bg) !important;
        color: var(--text) !important;
    }

    div[data-baseweb="input"] > div:focus-within,
    div[data-baseweb="textarea"] > div:focus-within,
    div[data-baseweb="select"] > div:focus-within {
        border-color: #FFFFFF !important;
        background-color: var(--field-bg-focus) !important;
        box-shadow: 0 0 0 1px #FFFFFF22 !important;
    }

    .stTextInput input,
    .stTextArea textarea,
    div[data-baseweb="select"] span,
    div[data-baseweb="select"] input {
        color: var(--text) !important;
        caret-color: var(--text) !important;
    }

    .stTextInput input::placeholder,
    .stTextArea textarea::placeholder {
        color: #9CA3AF !important;
    }

    /* 컬러피커/파일업로더/데이터에디터 보정 */
    [data-testid="stFileUploader"] section {
        background-color: var(--field-bg) !important;
        border: 1px dashed var(--line-strong) !important;
        border-radius: 14px !important;
        color: var(--text) !important;
    }

    [data-testid="stFileUploader"] section * {
        color: var(--text) !important;
    }

    [data-testid="stDataFrame"],
    [data-testid="stDataFrame"] div {
        border-color: var(--line) !important;
    }

    /* 버튼: 흰색 배경 + 검정 글씨 */
    div.stButton > button:first-child,
    .stDownloadButton > button {
        background-color: var(--button-bg);
        color: var(--button-text);
        border-radius: 12px;
        width: 100%;
        height: 44px;
        font-weight: 850;
        border: 1px solid var(--button-bg);
        box-shadow: none;
    }

    div.stButton > button:first-child p,
    .stDownloadButton > button p {
        color: var(--button-text) !important;
    }

    div.stButton > button:hover,
    .stDownloadButton > button:hover {
        background-color: #E5E5E5;
        color: #111111;
        border-color: #E5E5E5;
    }

    div.stButton > button:active,
    .stDownloadButton > button:active {
        background-color: #D4D4D4;
        color: #111111;
        border-color: #D4D4D4;
    }

    .naver-map-button {
        display: block;
        width: 100%;
        box-sizing: border-box;
        height: 44px;
        line-height: 42px;
        text-align: center;
        text-decoration: none !important;
        border-radius: 12px;
        border: 1px solid var(--button-bg);
        background: var(--button-bg);
        color: var(--button-text) !important;
        font-size: 14px;
        font-weight: 850;
        margin: 6px 0 8px 0;
    }

    .naver-map-button:hover {
        background: #E5E5E5;
        border-color: #E5E5E5;
        color: #111111 !important;
    }

    .capture-guide {
        margin: 6px 0 8px 0;
        padding: 10px 12px;
        border: 1px solid var(--line);
        border-radius: 12px;
        background: #111111;
        color: var(--muted);
        font-size: 12px;
        line-height: 18px;
    }

    .capture-guide strong {
        color: var(--text);
    }

    .mini-loading {
        display: flex;
        align-items: center;
        gap: 8px;
        margin: 8px 0 10px 0;
        padding: 9px 11px;
        border: 1px solid #3A3A3A;
        border-radius: 12px;
        background: #141414;
        color: #D6D6D6;
        font-size: 12px;
        line-height: 18px;
    }

    .mini-spinner {
        width: 14px;
        height: 14px;
        border: 2px solid #565656;
        border-top-color: #FFFFFF;
        border-radius: 50%;
        display: inline-block;
        animation: mini-spin 0.8s linear infinite;
        flex: 0 0 auto;
    }

    @keyframes mini-spin {
        to { transform: rotate(360deg); }
    }

    /* 탭 */
    [data-testid="stTabs"] button {
        color: var(--muted) !important;
    }

    [data-testid="stTabs"] button[aria-selected="true"] {
        color: var(--text) !important;
        font-weight: 850;
    }

    [data-testid="stTabs"] [data-baseweb="tab-highlight"] {
        background-color: #FFFFFF !important;
    }

    /* 미리보기 iframe 배경: 붉은 기 없는 쿨 그레이 */
    [data-testid="stTabs"] iframe {
        background-color: #EEF2F6 !important;
        border-radius: 14px;
    }

    hr {
        border-color: var(--line) !important;
    }

    .mini-help {
        color: var(--muted);
        font-size: 12px;
        line-height: 18px;
        margin-top: -4px;
    }


    /* 우측 미리보기 영역 고정 */
    div[data-testid="stHorizontalBlock"] > div[data-testid="column"]:nth-of-type(2) {
        position: sticky;
        top: 18px;
        align-self: flex-start;
        max-height: calc(100vh - 28px);
        overflow-y: auto;
        padding-right: 4px;
    }

    div[data-testid="stHorizontalBlock"] > div[data-testid="column"]:nth-of-type(2)::-webkit-scrollbar {
        width: 8px;
    }

    div[data-testid="stHorizontalBlock"] > div[data-testid="column"]:nth-of-type(2)::-webkit-scrollbar-thumb {
        background: #4A4A4A;
        border-radius: 999px;
    }

    div[data-testid="stHorizontalBlock"] > div[data-testid="column"]:nth-of-type(2)::-webkit-scrollbar-track {
        background: #111111;
    }

    /* 좌측 입력 설정 스크롤 패널 보정 */
    div[style*="overflow"]::-webkit-scrollbar {
        width: 8px;
        height: 8px;
    }

    div[style*="overflow"]::-webkit-scrollbar-thumb {
        background: #4A4A4A;
        border-radius: 999px;
    }

    div[style*="overflow"]::-webkit-scrollbar-track {
        background: #111111;
        border-radius: 999px;
    }



    .naver-local-card {
        border: 1px solid #3A3A3A;
        background: #242424;
        border-radius: 14px;
        padding: 12px 14px;
        margin: 8px 0 6px 0;
    }

    .naver-local-card-title {
        color: #FFFFFF;
        font-weight: 850;
        font-size: 14px;
        line-height: 1.35;
        margin-bottom: 5px;
    }

    .naver-local-card-meta {
        color: #B8B8B8;
        font-size: 12px;
        line-height: 1.55;
    }

    .naver-local-card-address {
        color: #F0F0F0;
        font-size: 12.5px;
        line-height: 1.55;
        margin-top: 4px;
    }

    /* 좌측 입력부 내부 카드 간격 압축 */
    div[data-testid="column"]:first-of-type [data-testid="stVerticalBlockBorderWrapper"] {
        margin-bottom: 10px;
    }
    </style>
    """,
    unsafe_allow_html=True,
)

if app_theme == "화이트 모드":
    st.markdown(
        """
        <style>
        :root {
            --app-bg: #F4F6F8;
            --card-bg: #FFFFFF;
            --card-bg-strong: #F8FAFC;
            --field-bg: #FFFFFF;
            --field-bg-focus: #F9FAFB;
            --line: #E5E7EB;
            --line-strong: #c9c9c9;
            --text: #111827;
            --muted: #64748B;
            --soft: #F1F5F9;
            --button-bg: #7b7b7b;
            --button-text: #FFFFFF;
        }

        [data-testid="stSidebar"] {
            background-color: #FFFFFF !important;
        }

        [data-testid="stVerticalBlockBorderWrapper"] {
            box-shadow: 0 14px 34px rgba(15, 23, 42, 0.08);
        }

        div[data-baseweb="input"] > div,
        div[data-baseweb="textarea"] > div,
        div[data-baseweb="select"] > div {
            background-color: #FFFFFF !important;
            border-color: #c9c9c9 !important;
        }

        input,
        textarea,
        [data-baseweb="input"],
        [data-baseweb="textarea"],
        [data-baseweb="select"] {
            border-color: #c9c9c9 !important;
        }

        [data-testid="stColorPicker"] input,
        [data-testid="stNumberInput"] input {
            border-color: #c9c9c9 !important;
        }

        div[data-baseweb="input"] > div:focus-within,
        div[data-baseweb="textarea"] > div:focus-within,
        div[data-baseweb="select"] > div:focus-within {
            border-color: #111827 !important;
            background-color: #FFFFFF !important;
            box-shadow: 0 0 0 1px rgba(17, 24, 39, 0.16) !important;
        }

        .stTextInput input,
        .stTextArea textarea,
        div[data-baseweb="select"] span,
        div[data-baseweb="select"] input {
            color: #111827 !important;
            caret-color: #111827 !important;
        }

        .stTextInput input::placeholder,
        .stTextArea textarea::placeholder {
            color: #94A3B8 !important;
        }

        [data-testid="stFileUploader"] section {
            background-color: #FFFFFF !important;
            border: 1px dashed #c9c9c9 !important;
            color: #111827 !important;
        }

        [data-testid="stFileUploader"] section * {
            color: #111827 !important;
        }

        div.stButton > button:first-child,
        .stDownloadButton > button,
        .naver-map-button,
        [data-testid="stFileUploader"] section button,
        [data-testid="stFileUploader"] button {
            border: 1px solid #c9c9c9 !important;
            border-radius: 12px !important;
            height: 40px !important;
            font-size: 13px !important;
            font-weight: 800 !important;
            color: #FFFFFF !important;
            background: #7b7b7b !important;
            cursor: pointer !important;
            box-shadow: none !important;
        }

        div.stButton > button:first-child p,
        .stDownloadButton > button p,
        [data-testid="stFileUploader"] section button p,
        [data-testid="stFileUploader"] button p {
            color: #FFFFFF !important;
        }

        div.stButton > button:hover,
        div.stButton > button:active,
        div.stButton > button:focus,
        .stDownloadButton > button:hover,
        .stDownloadButton > button:active,
        .stDownloadButton > button:focus,
        .naver-map-button:hover,
        .naver-map-button:active,
        .naver-map-button:focus,
        [data-testid="stFileUploader"] section button:hover,
        [data-testid="stFileUploader"] section button:active,
        [data-testid="stFileUploader"] section button:focus,
        [data-testid="stFileUploader"] button:hover,
        [data-testid="stFileUploader"] button:active,
        [data-testid="stFileUploader"] button:focus {
            background: #7b7b7b !important;
            color: #FFFFFF !important;
            border-color: #c9c9c9 !important;
            box-shadow: none !important;
        }

        .capture-guide,
        .mini-loading {
            background: #7b7b7b !important;
            border-color: #c9c9c9 !important;
            color: #FFFFFF !important;
        }

        .capture-guide strong {
            color: #FFFFFF !important;
        }

        .mini-spinner {
            border-color: #CBD5E1 !important;
            border-top-color: #111827 !important;
        }

        [data-testid="stTabs"] button {
            color: #64748B !important;
        }

        [data-testid="stTabs"] button[aria-selected="true"] {
            color: #111827 !important;
        }

        [data-testid="stTabs"] [data-baseweb="tab-highlight"] {
            background-color: #111827 !important;
        }

        .naver-local-card {
            border: 1px solid #E5E7EB !important;
            background: #FFFFFF !important;
            box-shadow: 0 8px 22px rgba(15, 23, 42, 0.06);
        }

        .naver-local-card-title {
            color: #111827 !important;
        }

        .naver-local-card-meta {
            color: #64748B !important;
        }

        .naver-local-card-address {
            color: #1F2937 !important;
        }

        .stApp,
        .stApp * {
            scrollbar-color: #7b7b7b #F1F3F5;
            scrollbar-width: thin;
        }

        .stApp *::-webkit-scrollbar,
        div[data-testid="stHorizontalBlock"] > div[data-testid="column"]:nth-of-type(2)::-webkit-scrollbar,
        div[style*="overflow"]::-webkit-scrollbar {
            width: 9px !important;
            height: 9px !important;
        }

        .stApp *::-webkit-scrollbar-thumb,
        div[data-testid="stHorizontalBlock"] > div[data-testid="column"]:nth-of-type(2)::-webkit-scrollbar-thumb,
        div[style*="overflow"]::-webkit-scrollbar-thumb {
            background: #7b7b7b !important;
            border-radius: 999px !important;
            border: 2px solid #F1F3F5 !important;
        }

        .stApp *::-webkit-scrollbar-thumb:hover,
        div[data-testid="stHorizontalBlock"] > div[data-testid="column"]:nth-of-type(2)::-webkit-scrollbar-thumb:hover,
        div[style*="overflow"]::-webkit-scrollbar-thumb:hover {
            background: #7b7b7b !important;
        }

        .stApp *::-webkit-scrollbar-track,
        div[data-testid="stHorizontalBlock"] > div[data-testid="column"]:nth-of-type(2)::-webkit-scrollbar-track,
        div[style*="overflow"]::-webkit-scrollbar-track {
            background: #F1F3F5 !important;
            border-radius: 999px !important;
        }

        /* 화이트 모드에서 남아 있는 검정 버튼/컨트롤류 보정 */
        button[kind="secondary"],
        button[kind="primary"],
        [role="button"]:not([aria-label*="Help"]),
        [data-testid="baseButton-secondary"],
        [data-testid="baseButton-primary"] {
            border: 1px solid #c9c9c9 !important;
            border-radius: 12px !important;
            font-size: 13px !important;
            font-weight: 800 !important;
            color: #FFFFFF !important;
            background: #7b7b7b !important;
            cursor: pointer !important;
            box-shadow: none !important;
        }

        button[kind="secondary"]:hover,
        button[kind="primary"]:hover,
        [role="button"]:not([aria-label*="Help"]):hover,
        [data-testid="baseButton-secondary"]:hover,
        [data-testid="baseButton-primary"]:hover {
            background: #7b7b7b !important;
            color: #FFFFFF !important;
            border-color: #c9c9c9 !important;
        }


        /* 화이트 모드: 도움말 툴팁/팝오버/드롭다운 메뉴 배경 보정 */
        [data-testid="stTooltipContent"],
        [data-testid="stTooltipContent"] *,
        [data-baseweb="tooltip"],
        [data-baseweb="tooltip"] *,
        [role="tooltip"],
        [role="tooltip"] *,
        div[data-baseweb="popover"],
        div[data-baseweb="popover"] *,
        div[data-baseweb="menu"],
        div[data-baseweb="menu"] *,
        ul[role="listbox"],
        ul[role="listbox"] *,
        div[role="listbox"],
        div[role="listbox"] * {
            background-color: #7b7b7b !important;
            color: #FFFFFF !important;
            border-color: #c9c9c9 !important;
        }

        [data-testid="stTooltipContent"],
        [data-baseweb="tooltip"],
        [role="tooltip"],
        div[data-baseweb="popover"],
        div[data-baseweb="menu"],
        ul[role="listbox"],
        div[role="listbox"] {
            border: 1px solid #c9c9c9 !important;
            border-radius: 12px !important;
            box-shadow: none !important;
            overflow: hidden !important;
        }

        li[role="option"],
        div[role="option"] {
            background: #7b7b7b !important;
            color: #FFFFFF !important;
            font-size: 13px !important;
            font-weight: 800 !important;
        }

        li[role="option"]:hover,
        div[role="option"]:hover,
        li[aria-selected="true"],
        div[aria-selected="true"] {
            background: #6f6f6f !important;
            color: #FFFFFF !important;
        }

        /* 화이트 모드: Data editor/DataFrame 영역 보정 */
        [data-testid="stDataFrame"],
        [data-testid="stDataFrame"] > div,
        [data-testid="stDataFrame"] div[role="grid"],
        [data-testid="stDataFrame"] .glideDataEditor {
            border: 1px solid #c9c9c9 !important;
            border-radius: 12px !important;
            background: #7b7b7b !important;
            color: #FFFFFF !important;
            box-shadow: none !important;
            overflow: hidden !important;
        }

        [data-testid="stDataFrame"] *,
        [data-testid="stDataFrame"] div,
        [data-testid="stDataFrame"] span,
        [data-testid="stDataFrame"] p,
        [data-testid="stDataFrame"] label {
            border-color: #c9c9c9 !important;
            color: #FFFFFF !important;
        }

        [data-testid="stDataFrame"] button,
        [data-testid="stDataFrame"] [role="button"],
        [data-testid="stDataFrame"] input,
        [data-testid="stDataFrame"] textarea {
            border: 1px solid #c9c9c9 !important;
            border-radius: 12px !important;
            min-height: 36px !important;
            font-size: 13px !important;
            font-weight: 800 !important;
            color: #FFFFFF !important;
            background: #7b7b7b !important;
            cursor: pointer !important;
            box-shadow: none !important;
        }

        [data-testid="stDataFrame"] button:hover,
        [data-testid="stDataFrame"] [role="button"]:hover,
        [data-testid="stDataFrame"] input:focus,
        [data-testid="stDataFrame"] textarea:focus {
            background: #7b7b7b !important;
            color: #FFFFFF !important;
            border-color: #c9c9c9 !important;
            box-shadow: none !important;
        }

        [data-testid="stDataFrame"] canvas {
            background: #7b7b7b !important;
            border-radius: 12px !important;
        }

        /* 화이트 모드: 검정 배경으로 남는 컨트롤류를 회색 톤으로 통일 */
        .stApp div[style*="background: #111111"],
        .stApp div[style*="background-color: #111111"],
        .stApp div[style*="background: rgb(17, 17, 17)"],
        .stApp div[style*="background-color: rgb(17, 17, 17)"],
        .stApp div[style*="background: #141414"],
        .stApp div[style*="background-color: #141414"],
        .stApp div[style*="background: rgb(20, 20, 20)"],
        .stApp div[style*="background-color: rgb(20, 20, 20)"] {
            background: #7b7b7b !important;
            background-color: #7b7b7b !important;
            border-color: #c9c9c9 !important;
            color: #FFFFFF !important;
        }

        .stApp div[style*="background: #111111"] *,
        .stApp div[style*="background-color: #111111"] *,
        .stApp div[style*="background: rgb(17, 17, 17)"] *,
        .stApp div[style*="background-color: rgb(17, 17, 17)"] *,
        .stApp div[style*="background: #141414"] *,
        .stApp div[style*="background-color: #141414"] *,
        .stApp div[style*="background: rgb(20, 20, 20)"] *,
        .stApp div[style*="background-color: rgb(20, 20, 20)"] * {
            color: #FFFFFF !important;
            border-color: #c9c9c9 !important;
        }

        hr {
            border-color: #E5E7EB !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )



# 화이트 모드: Glide Data Editor 내부 검정 배경 추가 보정
# st.data_editor는 내부 일부를 canvas/Glide Data Grid로 렌더링하기 때문에
# 일반 stDataFrame 선택자만으로는 검정 배경이 남을 수 있어 전용 클래스를 한 번 더 보정합니다.
if app_theme == "화이트 모드":
    st.markdown(
        """
        <style>
        /* 커리큘럼/문의처 data_editor 최외곽 */
        .stApp .st-key-curr_editor,
        .stApp .st-key-contacts_editor,
        .stApp .st-key-curr_editor [data-testid="stDataFrame"],
        .stApp .st-key-contacts_editor [data-testid="stDataFrame"] {
            border: 1px solid #c9c9c9 !important;
            border-radius: 12px !important;
            background: #7b7b7b !important;
            color: #FFFFFF !important;
            box-shadow: none !important;
            overflow: hidden !important;
        }

        /* Glide Data Editor 내부 스크롤 영역 */
        .stApp .st-key-curr_editor .stDataFrameGlideDataEditor,
        .stApp .st-key-contacts_editor .stDataFrameGlideDataEditor,
        .stApp .st-key-curr_editor .dvn-scroller,
        .stApp .st-key-contacts_editor .dvn-scroller,
        .stApp div.stDataFrameGlideDataEditor,
        .stApp div.dvn-scroller.stDataFrameGlideDataEditor {
            border: 1px solid #c9c9c9 !important;
            border-radius: 12px !important;
            background: #7b7b7b !important;
            background-color: #7b7b7b !important;
            color: #FFFFFF !important;
            box-shadow: none !important;
            overflow: hidden !important;
        }

        /* Glide 내부 wrapper까지 검정 배경 제거 */
        .stApp .stDataFrameGlideDataEditor > div,
        .stApp .stDataFrameGlideDataEditor > div > div,
        .stApp .dvn-scroller > div,
        .stApp .dvn-scroller > div > div {
            background: #7b7b7b !important;
            background-color: #7b7b7b !important;
            border-color: #c9c9c9 !important;
            color: #FFFFFF !important;
        }

        /* canvas 자체의 배경 보정 */
        .stApp .stDataFrameGlideDataEditor canvas,
        .stApp .dvn-scroller canvas,
        .stApp [data-testid="stDataFrame"] canvas {
            background: #7b7b7b !important;
            background-color: #7b7b7b !important;
            border-radius: 12px !important;
        }

        /* data_editor 내부 툴바/셀렉터/버튼 계열 */
        .stApp .stDataFrameGlideDataEditor button,
        .stApp .stDataFrameGlideDataEditor [role="button"],
        .stApp .stDataFrameGlideDataEditor input,
        .stApp .stDataFrameGlideDataEditor textarea,
        .stApp .dvn-scroller button,
        .stApp .dvn-scroller [role="button"],
        .stApp .dvn-scroller input,
        .stApp .dvn-scroller textarea {
            border: 1px solid #c9c9c9 !important;
            border-radius: 12px !important;
            min-height: 36px !important;
            font-size: 13px !important;
            font-weight: 800 !important;
            color: #FFFFFF !important;
            background: #7b7b7b !important;
            background-color: #7b7b7b !important;
            cursor: pointer !important;
            box-shadow: none !important;
        }

        /* 내부 스크롤바 */
        .stApp .stDataFrameGlideDataEditor::-webkit-scrollbar,
        .stApp .dvn-scroller::-webkit-scrollbar,
        .stApp .stDataFrameGlideDataEditor *::-webkit-scrollbar,
        .stApp .dvn-scroller *::-webkit-scrollbar {
            width: 9px !important;
            height: 9px !important;
        }

        .stApp .stDataFrameGlideDataEditor::-webkit-scrollbar-thumb,
        .stApp .dvn-scroller::-webkit-scrollbar-thumb,
        .stApp .stDataFrameGlideDataEditor *::-webkit-scrollbar-thumb,
        .stApp .dvn-scroller *::-webkit-scrollbar-thumb {
            background: #7b7b7b !important;
            border: 2px solid #F1F3F5 !important;
            border-radius: 999px !important;
        }

        .stApp .stDataFrameGlideDataEditor::-webkit-scrollbar-track,
        .stApp .dvn-scroller::-webkit-scrollbar-track,
        .stApp .stDataFrameGlideDataEditor *::-webkit-scrollbar-track,
        .stApp .dvn-scroller *::-webkit-scrollbar-track {
            background: #F1F3F5 !important;
            border-radius: 999px !important;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

top_title_col, top_theme_col = st.columns([1.0, 0.24], vertical_alignment="top")
with top_title_col:
    st.markdown('<div class="main-title">✉️ 교육 안내문 작성 도구 (Beta)</div>', unsafe_allow_html=True)
    st.markdown(
        '<div class="sub-title">교육 안내문 작성과 아웃룩 본문용 HTML·PNG/JPG 생성을 지원하는 내부 테스트용 도구입니다.</div>',
        unsafe_allow_html=True,
    )
with top_theme_col:
    st.markdown('<div class="theme-select-note">화면 테마</div>', unsafe_allow_html=True)
    st.selectbox(
        "화면 테마",
        ["다크 모드", "화이트 모드"],
        key="app_theme",
        label_visibility="collapsed",
        help="작업 화면 테마만 변경합니다. 오른쪽 안내문 결과물은 메일 발송용 흰 배경으로 유지됩니다.",
    )


# -----------------------------
# 지도 이미지 클립보드 붙여넣기 컴포넌트
# -----------------------------
def clipboard_map_paste_component() -> None:
    is_light_theme = st.session_state.get("app_theme") == "화이트 모드"
    paste_box_bg = "#7b7b7b" if is_light_theme else "#111111"
    paste_box_border = "#c9c9c9" if is_light_theme else "#7A7A7A"
    paste_box_focus = "#c9c9c9" if is_light_theme else "#FFFFFF"
    paste_box_focus_shadow = "#c9c9c955" if is_light_theme else "#FFFFFF33"
    paste_box_title = "#FFFFFF" if is_light_theme else "#FFFFFF"
    paste_box_desc = "#F2F2F2" if is_light_theme else "#B8B8B8"
    paste_box_status = "#F5F5F5" if is_light_theme else "#D6D6D6"
    components.html(
        """
<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8" />
<style>
    body { margin: 0; padding: 0; background: transparent; font-family: Arial, sans-serif; }
    .paste-box { box-sizing: border-box; width: 100%; min-height: 112px; border: 1px dashed __PASTE_BOX_BORDER__; border-radius: 14px; background: __PASTE_BOX_BG__; color: __PASTE_BOX_TITLE__; display: flex; flex-direction: column; align-items: center; justify-content: center; gap: 7px; cursor: text; text-align: center; outline: none; padding: 14px 16px; }
    .paste-box:focus { border-color: __PASTE_BOX_FOCUS__; box-shadow: 0 0 0 1px __PASTE_BOX_FOCUS_SHADOW__; }
    .title { font-size: 14px; font-weight: 850; color: __PASTE_BOX_TITLE__; }
    .desc { font-size: 12px; line-height: 18px; color: __PASTE_BOX_DESC__; }
    .status { font-size: 12px; line-height: 18px; color: __PASTE_BOX_STATUS__; min-height: 18px; }
</style>
</head>
<body>
    <div id="pasteBox" class="paste-box" tabindex="0">
        <div class="title">지도 이미지 붙여넣기</div>
        <div class="desc">네이버 지도에서 원하는 영역을 캡처한 뒤<br>이 박스를 클릭하고 Ctrl + V를 누르세요.</div>
        <div id="status" class="status">파일 저장 없이 바로 적용할 수 있습니다.</div>
    </div>
<script>
const pasteBox = document.getElementById('pasteBox');
const statusEl = document.getElementById('status');
function setStatus(text) { statusEl.textContent = text; }
function getParentDocument() { try { return window.parent.document; } catch (e) { return null; } }
function setNativeTextareaValue(textarea, value) {
    const win = textarea.ownerDocument.defaultView || window.parent || window;
    const descriptor = Object.getOwnPropertyDescriptor(win.HTMLTextAreaElement.prototype, 'value');
    if (descriptor && descriptor.set) { descriptor.set.call(textarea, value); } else { textarea.value = value; }
    textarea.focus();
    textarea.dispatchEvent(new win.InputEvent('input', { bubbles: true, inputType: 'insertText', data: value }));
    textarea.dispatchEvent(new win.Event('change', { bubbles: true }));
    textarea.blur();
}
function clickRefreshButton(doc) {
    const buttons = Array.from(doc.querySelectorAll('button'));
    const button = buttons.find((btn) => (btn.innerText || btn.textContent || '').indexOf('붙여넣기 반영') >= 0);
    if (button) {
        setTimeout(() => button.click(), 250);
        return true;
    }
    return false;
}
function findTargetTextarea(doc) {
    const textareas = Array.from(doc.querySelectorAll('textarea'));
    return textareas.find((el) => (el.getAttribute('aria-label') || '').indexOf('클립보드 이미지 데이터') >= 0) || textareas.find((el) => (el.placeholder || '').indexOf('클립보드') >= 0) || null;
}
function applyDataUrlToStreamlit(dataUrl) {
    const doc = getParentDocument();
    if (!doc) return false;
    const textarea = findTargetTextarea(doc);
    if (!textarea) return false;
    try {
        setNativeTextareaValue(textarea, dataUrl);
        clickRefreshButton(doc);
        return true;
    } catch (e) { console.log(e); return false; }
}
function imageFileToDataUrl(file) {
    return new Promise((resolve, reject) => {
        const img = new Image();
        img.onload = function() {
            const maxWidth = 1200;
            const scale = Math.min(1, maxWidth / img.width);
            const width = Math.max(1, Math.round(img.width * scale));
            const height = Math.max(1, Math.round(img.height * scale));
            const canvas = document.createElement('canvas');
            canvas.width = width; canvas.height = height;
            const ctx = canvas.getContext('2d');
            ctx.drawImage(img, 0, 0, width, height);
            const dataUrl = canvas.toDataURL('image/jpeg', 0.9);
            URL.revokeObjectURL(img.src);
            resolve(dataUrl);
        };
        img.onerror = reject;
        img.src = URL.createObjectURL(file);
    });
}
async function handlePaste(event) {
    const items = event.clipboardData && event.clipboardData.items ? Array.from(event.clipboardData.items) : [];
    const imageItem = items.find((item) => item.type && item.type.indexOf('image/') === 0);
    if (!imageItem) { setStatus('붙여넣은 내용에서 이미지를 찾지 못했습니다.'); return; }
    event.preventDefault();
    const file = imageItem.getAsFile();
    if (!file) { setStatus('이미지를 읽지 못했습니다.'); return; }
    setStatus('이미지를 적용하는 중입니다...');
    try {
        const dataUrl = await imageFileToDataUrl(file);
        const ok = applyDataUrlToStreamlit(dataUrl);
        setStatus(ok ? '지도 이미지를 붙여넣었습니다. 화면을 갱신하는 중입니다.' : '자동 적용 실패: 파일 첨부를 사용해 주세요.');
    } catch (e) { console.log(e); setStatus('이미지 처리 중 오류가 발생했습니다. 파일 첨부를 사용해 주세요.'); }
}
pasteBox.addEventListener('click', () => pasteBox.focus());
pasteBox.addEventListener('paste', handlePaste);
document.addEventListener('paste', function(event) { if (document.activeElement === pasteBox) handlePaste(event); });
</script>
</body>
</html>
        """
        .replace("__PASTE_BOX_BG__", paste_box_bg)
        .replace("__PASTE_BOX_BORDER__", paste_box_border)
        .replace("__PASTE_BOX_FOCUS__", paste_box_focus)
        .replace("__PASTE_BOX_FOCUS_SHADOW__", paste_box_focus_shadow)
        .replace("__PASTE_BOX_TITLE__", paste_box_title)
        .replace("__PASTE_BOX_DESC__", paste_box_desc)
        .replace("__PASTE_BOX_STATUS__", paste_box_status),
        height=126,
    )


# -----------------------------
# 스포이드 컴포넌트
# -----------------------------
def eyedropper_component(target: str, label: str) -> None:
    safe_target = esc_attr(target)
    safe_label = esc(label)
    target_label = "메인 컬러" if target == "main" else "하단 컬러"
    components.html(
        f"""
<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8" />
<style>
    body {{ margin: 0; padding: 0; background: transparent; font-family: Arial, sans-serif; }}
    .wrap {{ display: flex; align-items: flex-start; gap: 6px; min-height: 64px; flex-wrap: wrap; padding: 2px 0; box-sizing: border-box; overflow: visible; }}
    button {{
        height: 32px;
        padding: 0 10px;
        border-radius: 9px;
        border: 1px solid #FFFFFF;
        background: #FFFFFF;
        color: #111111 !important;
        font-size: 11px;
        font-weight: 800;
        cursor: pointer;
        text-decoration: none;
        display: inline-flex;
        align-items: center;
        justify-content: center;
        box-sizing: border-box;
    }}
    button:hover {{ background: #E5E5E5; }}
    .result {{ font-size: 12px; color: #B8B8B8; white-space: nowrap; line-height: 32px; }}
    #fallback button {{ margin-top: 0; }}
</style>
</head>
<body>
    <div class="wrap">
        <button onclick="pickColor()">🎯 {safe_label}</button>
        <span id="result" class="result">웹 화면 색상 추출</span>
        <span id="fallback"></span>
    </div>
<script>
const TARGET = '{safe_target}';
const TARGET_LABEL = '{target_label}';

function getParentDocument() {{
    try {{ return window.parent.document; }} catch (e) {{ return null; }}
}}

function setNativeInputValue(input, value) {{
    const win = input.ownerDocument.defaultView || window.parent || window;
    const descriptor = Object.getOwnPropertyDescriptor(win.HTMLInputElement.prototype, 'value');
    if (descriptor && descriptor.set) {{
        descriptor.set.call(input, value);
    }} else {{
        input.value = value;
    }}
    input.dispatchEvent(new win.Event('input', {{ bubbles: true }}));
    input.dispatchEvent(new win.Event('change', {{ bubbles: true }}));
}}

function findColorPickerInput(doc) {{
    const pickers = Array.from(doc.querySelectorAll('[data-testid="stColorPicker"]'));
    const matchedPicker = pickers.find((picker) => {{
        const txt = picker.innerText || picker.textContent || '';
        return txt.indexOf(TARGET_LABEL) >= 0;
    }});

    if (matchedPicker) {{
        const matchedInput = matchedPicker.querySelector('input[id^="rc-editable-input-"], input[type="text"], input');
        if (matchedInput) return matchedInput;
    }}

    const visibleInputs = Array.from(doc.querySelectorAll('input[id^="rc-editable-input-"]'))
        .filter((input) => input.offsetParent !== null);
    if (TARGET === 'footer') return visibleInputs[1] || visibleInputs[0] || null;
    return visibleInputs[0] || null;
}}

function applyColorDirectlyToStreamlit(color) {{
    const doc = getParentDocument();
    if (!doc) return false;

    const input = findColorPickerInput(doc);
    if (!input) return false;

    try {{
        input.scrollIntoView({{ block: 'center', inline: 'nearest' }});
        input.focus();
        input.select && input.select();
        setNativeInputValue(input, color);

        const win = input.ownerDocument.defaultView || window.parent || window;
        input.dispatchEvent(new win.KeyboardEvent('keydown', {{ key: 'Enter', code: 'Enter', bubbles: true }}));
        input.dispatchEvent(new win.KeyboardEvent('keyup', {{ key: 'Enter', code: 'Enter', bubbles: true }}));
        input.blur();
        return true;
    }} catch (e) {{
        console.log(e);
        return false;
    }}
}}

function getStreamlitUrl() {{
    let baseUrl = document.referrer || '';
    try {{
        if (!baseUrl) baseUrl = window.top.location.href;
    }} catch (e) {{}}
    if (!baseUrl) baseUrl = window.location.href;
    return new URL(baseUrl);
}}

function buildApplyUrl(color) {{
    const url = getStreamlitUrl();
    url.searchParams.set('picked_target', TARGET);
    url.searchParams.set('picked_color', color);
    url.searchParams.set('picked_token', String(Date.now()));
    url.hash = '';
    return url.toString();
}}

function copyTextFallback(text) {{
    try {{
        const textarea = document.createElement('textarea');
        textarea.value = text;
        textarea.setAttribute('readonly', '');
        textarea.style.position = 'fixed';
        textarea.style.left = '-9999px';
        document.body.appendChild(textarea);
        textarea.select();
        document.execCommand('copy');
        document.body.removeChild(textarea);
        return true;
    }} catch (e) {{
        return false;
    }}
}}

async function copyColorCode(color) {{
    let ok = false;
    try {{
        await navigator.clipboard.writeText(color);
        ok = true;
    }} catch (e) {{
        ok = copyTextFallback(color);
    }}

    const resultEl = document.getElementById('result');
    if (ok) {{
        resultEl.textContent = color + ' 복사 완료';
    }} else {{
        resultEl.textContent = color + ' 복사 실패';
        alert('자동 복사가 막혔습니다. 색상 코드를 직접 복사해 주세요: ' + color);
    }}
}}

function showCopyButton(color) {{
    const fallbackEl = document.getElementById('fallback');
    fallbackEl.innerHTML = '';
    const button = document.createElement('button');
    button.type = 'button';
    button.textContent = '색상 코드 복사하기';
    button.onclick = function() {{ copyColorCode(color); }};
    fallbackEl.appendChild(button);
}}

async function pickColor() {{
    const resultEl = document.getElementById('result');
    if (!window.EyeDropper) {{
        resultEl.textContent = '현재 브라우저 미지원';
        alert('이 브라우저에서는 스포이드 기능을 지원하지 않습니다. Chrome 또는 Edge에서 실행해 주세요.');
        return;
    }}

    try {{
        const eyeDropper = new EyeDropper();
        const result = await eyeDropper.open();
        const color = result.sRGBHex;

        if (applyColorDirectlyToStreamlit(color)) {{
            resultEl.textContent = color + ' 입력창 적용 완료';
            return;
        }}

        resultEl.textContent = color + ' 추출 완료';
        showCopyButton(color);
    }} catch (error) {{
        resultEl.textContent = '취소됨';
    }}
}}
</script>
</body>
</html>
        """,
        height=92,
    )


# -----------------------------
# 안내문 HTML 빌더
# -----------------------------
def build_section_title(number: int, title: str, main_color: str, title_text_color: str = "#FFFFFF", zone: str = "") -> str:
    return f"""
    <table role="presentation" cellspacing="0" cellpadding="0" border="0" data-zone="{esc_attr(zone)} section-title brand header_text" style="margin: 0 0 14px 0; border-collapse: collapse;">
        <tr>
            <td style="width: 6px; background-color: {main_color}; font-size: 0; line-height: 0;">&nbsp;</td>
            <td style="background-color: {main_color}; color: {title_text_color}; padding: 8px 18px 8px 12px; font-size: 18px; line-height: 22px; font-weight: 700; letter-spacing: -0.3px;">
                {number}. {esc(title)}
            </td>
        </tr>
    </table>
    """


def build_bullet_rows(lines: list[str], text_color: str = "#343A40") -> str:
    rows = ""
    for line in lines:
        clean_line = normalize_bullet(line)
        if not clean_line:
            continue
        rows += f"""
        <tr>
            <td valign="top" style="width: 18px; padding: 0 8px 8px 0; font-size: 15px; color: {text_color}; line-height: 22px;">-</td>
            <td valign="top" style="padding: 0 0 8px 0; font-size: 15px; color: {text_color}; line-height: 22px; letter-spacing: -0.2px;">{esc(clean_line)}</td>
        </tr>
        """
    return rows or """
        <tr>
            <td style="font-size: 15px; color: #6B7280; line-height: 22px;">입력된 안내 사항이 없습니다.</td>
        </tr>
    """


def build_contact_rows(contacts: list[dict], text_color: str = "#343A40") -> str:
    rows = ""
    for contact in to_records(contacts):
        role = str(contact.get("role", "") or "").strip()
        phone = str(contact.get("phone", "") or "").strip()
        if not role and not phone:
            continue

        phone_html = f" <span style='color:{text_color};'>({esc(phone)})</span>" if phone else ""
        rows += f"""
        <tr>
            <td valign="top" style="width: 18px; padding: 0 8px 8px 0; font-size: 15px; color: {text_color}; line-height: 22px;">-</td>
            <td valign="top" style="padding: 0 0 8px 0; font-size: 15px; color: {text_color}; line-height: 22px;">
                <span style="font-weight: 700;">{esc(role)}</span>{phone_html}
            </td>
        </tr>
        """

    return rows or """
        <tr>
            <td style="font-size: 15px; color: #6B7280; line-height: 22px;">입력된 문의처가 없습니다.</td>
        </tr>
    """


def parse_curriculum_columns(columns_text: str) -> list[str]:
    """쉼표 또는 줄바꿈으로 입력한 커리큘럼 열 이름을 정리합니다."""
    raw = str(columns_text or "").replace("\n", ",")
    columns = []
    for part in raw.split(","):
        col = part.strip()
        if col and col not in columns:
            columns.append(col)
    return columns or ["Day", "시간", "교육 내용", "강사/비고"]




def parse_curriculum_column_defs(column_defs: object) -> list[str]:
    """data_editor에서 편집한 열 이름 목록을 실제 커리큘럼 컬럼명으로 변환합니다."""
    columns = []
    for row in to_records(column_defs):
        col = str(row.get("column_name", "") or row.get("열 이름", "") or row.get("name", "")).strip()
        if col and col not in columns:
            columns.append(col)
    return columns or ["Day", "시간", "교육 내용", "강사/비고"]


def _cell_value_by_column(row: dict, column: str) -> str:
    aliases = {
        "시간": ["시간", "time"],
        "Day": ["Day", "day", "일차"],
        "일차": ["일차", "day", "Day"],
        "교육 내용": ["교육 내용", "과정 내용", "subject"],
        "강사/비고": ["강사/비고", "비고", "강사", "speaker"],
    }
    keys = aliases.get(column, [column])
    for key in keys:
        if key in row and str(row.get(key, "") or "").strip():
            return str(row.get(key, "") or "").strip()
    return str(row.get(column, "") or "").strip()


def normalize_curriculum_for_columns(curriculum: list[dict], columns: list[str]) -> list[dict]:
    records = []
    for row in to_records(curriculum):
        normalized = {column: _cell_value_by_column(row, column) for column in columns}
        if any(str(value or "").strip() for value in normalized.values()):
            records.append(normalized)
    if not records:
        records = [{column: "" for column in columns}]
    return records



def _normalize_excel_header(value: object) -> str:
    text = str(value or "").strip().lower()
    text = re.sub(r"\s+", "", text)
    text = text.replace("_", "").replace("-", "").replace("/", "")
    return text


def _excel_cell_to_text(value: object) -> str:
    """엑셀 셀 값을 시간표 입력용 문자열로 변환합니다."""
    if value is None:
        return ""
    try:
        from datetime import date as _date, datetime as _datetime, time as _time
        if isinstance(value, _datetime):
            if value.time() == _time(0, 0):
                return value.strftime("%Y-%m-%d")
            return value.strftime("%Y-%m-%d %H:%M")
        if isinstance(value, _date):
            return value.strftime("%Y-%m-%d")
        if isinstance(value, _time):
            return value.strftime("%H:%M")
    except Exception:
        pass
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value).strip()


def parse_curriculum_excel(uploaded_file, target_columns: list[str]) -> tuple[list[dict], list[str], str]:
    """첫 번째 워크시트의 시간표를 읽어 현재 커리큘럼 열 구조에 맞춰 변환합니다.

    지원 형식은 .xlsx/.xlsm입니다. 첫 20행 안에서 Day·시간·교육 내용·강사/비고와
    유사한 헤더가 있는 행을 자동 탐지하며, 병합 셀 때문에 비어 있는 Day 값은 아래 행으로 이어받습니다.
    """
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        raise RuntimeError("엑셀 업로드에는 openpyxl이 필요합니다. `python -m pip install openpyxl`을 실행해 주세요.") from exc

    if uploaded_file is None:
        raise ValueError("엑셀 파일을 먼저 첨부해 주세요.")

    file_name = str(getattr(uploaded_file, "name", "") or "").lower()
    if not file_name.endswith((".xlsx", ".xlsm")):
        raise ValueError("현재는 .xlsx 또는 .xlsm 파일만 지원합니다.")

    workbook = load_workbook(io.BytesIO(uploaded_file.getvalue()), data_only=True, read_only=False)
    worksheet = workbook.active
    values = [list(row) for row in worksheet.iter_rows(values_only=True)]
    if not values:
        raise ValueError("엑셀 파일에서 데이터를 찾지 못했습니다.")

    alias_groups = {
        "Day": {"day", "일차", "일정", "일자", "차수"},
        "시간": {"시간", "교육시간", "강의시간", "time"},
        "교육 내용": {"교육내용", "과정내용", "강의내용", "내용", "주제", "모듈", "subject"},
        "강사/비고": {"강사비고", "강사", "비고", "담당", "speaker", "remark", "remarks"},
    }
    normalized_aliases = {
        canonical: {_normalize_excel_header(alias) for alias in aliases | {canonical}}
        for canonical, aliases in alias_groups.items()
    }

    best_index = -1
    best_score = 0
    for row_index, row in enumerate(values[:20]):
        normalized = [_normalize_excel_header(cell) for cell in row]
        score = sum(
            1
            for aliases in normalized_aliases.values()
            if any(cell and cell in aliases for cell in normalized)
        )
        if score > best_score:
            best_index = row_index
            best_score = score

    if best_index < 0 or best_score < 2:
        # 명확한 헤더가 없으면 첫 번째 비어 있지 않은 행을 헤더로 사용합니다.
        best_index = next((idx for idx, row in enumerate(values) if any(str(cell or "").strip() for cell in row)), -1)
        if best_index < 0:
            raise ValueError("엑셀 파일에서 헤더 행을 찾지 못했습니다.")

    raw_headers = [_excel_cell_to_text(cell) for cell in values[best_index]]
    normalized_headers = [_normalize_excel_header(cell) for cell in raw_headers]

    canonical_source_indexes: dict[str, int] = {}
    for canonical, aliases in normalized_aliases.items():
        for index, header in enumerate(normalized_headers):
            if header and header in aliases:
                canonical_source_indexes[canonical] = index
                break

    # 열명이 특이한 파일은 앞 4개 열을 Day/시간/교육 내용/강사·비고 순으로 보조 매핑합니다.
    for fallback_index, canonical in enumerate(["Day", "시간", "교육 내용", "강사/비고"]):
        if canonical not in canonical_source_indexes and fallback_index < len(raw_headers):
            canonical_source_indexes[canonical] = fallback_index

    output_columns = order_curriculum_columns(target_columns or ["Day", "시간", "교육 내용", "강사/비고"])
    rows: list[dict] = []
    previous_day = ""

    for raw_row in values[best_index + 1:]:
        canonical_values: dict[str, str] = {}
        for canonical, source_index in canonical_source_indexes.items():
            cell = raw_row[source_index] if source_index < len(raw_row) else None
            canonical_values[canonical] = _excel_cell_to_text(cell)

        if canonical_values.get("Day"):
            previous_day = canonical_values["Day"]
        elif previous_day:
            canonical_values["Day"] = previous_day

        mapped: dict[str, str] = {}
        for column in output_columns:
            normalized_column = _normalize_excel_header(column)
            if normalized_column in normalized_aliases["Day"]:
                mapped[column] = canonical_values.get("Day", "")
            elif normalized_column in normalized_aliases["시간"]:
                mapped[column] = canonical_values.get("시간", "")
            elif normalized_column in normalized_aliases["교육 내용"]:
                mapped[column] = canonical_values.get("교육 내용", "")
            elif normalized_column in normalized_aliases["강사/비고"]:
                mapped[column] = canonical_values.get("강사/비고", "")
            else:
                # 동일한 원본 헤더명이 있으면 사용자 정의 열에도 반영합니다.
                try:
                    source_index = normalized_headers.index(normalized_column)
                except ValueError:
                    source_index = -1
                mapped[column] = _excel_cell_to_text(raw_row[source_index]) if 0 <= source_index < len(raw_row) else ""

        if any(str(value or "").strip() for value in mapped.values()):
            rows.append(mapped)

    if not rows:
        raise ValueError("헤더 아래에서 시간표 데이터를 찾지 못했습니다.")

    return rows, output_columns, f"{worksheet.title} 시트에서 시간표 {len(rows)}행을 불러왔습니다."

def build_curriculum_header(columns: list[str], main_color: str, text_color: str) -> str:
    header_cells = ""
    total = len(columns)
    for idx, column in enumerate(columns):
        border = "border-right: 1px solid #74787C;" if idx < total - 1 else ""
        header_cells += f"""
                    <th style="padding: 12px 10px; background-color: {main_color}; color: {text_color}; font-size: 15px; line-height: 20px; font-weight: 700; text-align: center; {border}">{esc(column)}</th>"""
    return f"""
                <tr data-zone="curriculum brand header_text">{header_cells}
                </tr>"""


def build_curriculum_rows(
    curriculum: list[dict],
    columns: list[str] | None = None,
    odd_row_color: str = "#FFFFFF",
    even_row_color: str = "#FFFFFF",
) -> str:
    columns = columns or ["Day", "시간", "교육 내용", "강사/비고"]
    rows = ""
    normalized_records = normalize_curriculum_for_columns(curriculum, columns)

    for row_index, row in enumerate(normalized_records):
        if not any(str(row.get(column, "") or "").strip() for column in columns):
            continue
        row_background = odd_row_color if row_index % 2 == 0 else even_row_color
        cells = ""
        for idx, column in enumerate(columns):
            value = _cell_value_by_column(row, column)
            is_long_text = any(token in column for token in ["내용", "주제", "과정", "비고", "메모"])
            align = "left" if is_long_text else "center"
            weight = "700" if is_long_text else "500"
            border = "border-right: 1px solid #D8DEE6;" if idx < len(columns) - 1 else ""
            cells += f"""
            <td style="padding: 13px 10px; border-bottom: 1px solid #D8DEE6; {border} background-color: {row_background}; font-size: 14px; color: #1F2933; text-align: {align}; line-height: 20px; font-weight: {weight};">
                {esc(value) if value else "-"}
            </td>"""
        rows += f"""
        <tr>{cells}
        </tr>
        """

    return rows or f"""
        <tr>
            <td colspan="{len(columns)}" style="padding: 18px 10px; border-bottom: 1px solid #D8DEE6; background-color: #FFFFFF; font-size: 14px; color: #6B7280; text-align: center;">
                입력된 커리큘럼이 없습니다.
            </td>
        </tr>
    """


def build_logo_image_html(logo_image_data_url: str, logo_max_height: int = 52) -> str:
    if not str(logo_image_data_url or "").strip():
        return ""
    return f"""
    <img src="{esc_attr(logo_image_data_url)}" alt="회사 로고" style="display: inline-block; max-height: {int(logo_max_height)}px; max-width: 240px; width: auto; height: auto; border: 0; outline: none; text-decoration: none;" />
    """


def build_logo_row(logo_image_data_url: str, logo_position: str, area: str, logo_max_height: int = 52) -> str:
    if not str(logo_image_data_url or "").strip():
        return ""

    is_top = "상단" in logo_position
    is_bottom = "하단" in logo_position
    if area == "top" and not is_top:
        return ""
    if area == "bottom" and not is_bottom:
        return ""

    align = "right" if "우측" in logo_position else "left"
    margin = "0 0 12px 0" if area == "top" else "18px 0 0 0"
    return f"""
    <table role="presentation" cellspacing="0" cellpadding="0" border="0" data-zone="logo" style="width: 100%; border-collapse: collapse; margin: {margin};">
        <tr>
            <td style="text-align: {align}; vertical-align: middle; height: {int(logo_max_height)}px;">
                {build_logo_image_html(logo_image_data_url, logo_max_height)}
            </td>
        </tr>
    </table>
    """


def build_final_mail_html(
    company_name: str,
    course_name: str,
    date_range: str,
    day1_time: str,
    day2_time: str,
    place_name: str,
    road_address: str,
    map_image_src: str,
    delivery_type: str,
    welcome_title: str,
    welcome_body_text: str,
    time_notice_text: str,
    edited_curriculum: list[dict],
    info_text: str,
    edited_contacts: list[dict],
    main_color: str,
    footer_color: str,
    curriculum_columns_text: str = "Day, 시간, 교육 내용, 강사/비고",
    curriculum_title: str = "상세 커리큘럼",
    show_curriculum_table: bool = True,
    section_text_color: str = "#343A40",
    curr_header_text_color: str = "#FFFFFF",
    logo_image_data_url: str = "",
    logo_position: str = "우측 상단",
    logo_max_height: int = 52,
    font_stack: str = "'Malgun Gothic', 'Apple SD Gothic Neo', Arial, sans-serif",
    embedded_font_css: str = "",
    zoom_url: str = "",
    overview_section_title: str = "교육 개요",
    location_section_title: str = "교육 장소",
    notice_section_title: str = "안내 사항",
    contact_section_title: str = "관련 문의",
    curriculum_odd_row_color: str = "#FFFFFF",
    curriculum_even_row_color: str = "#FFFFFF",
) -> str:
    company_display = str(company_name or "").strip()
    full_course_name = f"{company_display} {course_name}".strip()
    location_lines = []
    for location_part in [place_name, road_address]:
        for location_line in str(location_part or "").splitlines():
            clean_location_line = location_line.strip()
            if clean_location_line and clean_location_line not in location_lines:
                location_lines.append(clean_location_line)
    curriculum_columns = parse_curriculum_columns(curriculum_columns_text)
    curriculum_header_html = build_curriculum_header(curriculum_columns, main_color, curr_header_text_color)
    curriculum_html = build_curriculum_rows(
        edited_curriculum,
        curriculum_columns,
        odd_row_color=curriculum_odd_row_color,
        even_row_color=curriculum_even_row_color,
    )
    curriculum_title_html = ""
    if str(curriculum_title or "").strip():
        curriculum_title_html = f"""
        <p style="margin: 4px 0 8px 0; font-size: 15px; line-height: 21px; color: #222222; font-weight: 800;">{esc(curriculum_title)}</p>
        """
    curriculum_block_html = ""
    if show_curriculum_table:
        curriculum_block_html = f"""
        {curriculum_title_html}
        <table role="presentation" cellspacing="0" cellpadding="0" border="0" data-zone="curriculum" style="width: 100%; border-collapse: collapse; margin: 10px 0 34px 0; border-top: 1px solid #D8DEE6; border-left: 1px solid #D8DEE6; background-color: #FFFFFF;">
            <thead>
                {curriculum_header_html}
            </thead>
            <tbody style="background-color: #FFFFFF;">
                {curriculum_html}
            </tbody>
        </table>
        """
    overview_rows = [
        f"교육명 : {full_course_name}",
        f"교육일정 : {date_range}  *집합 교육 기준",
    ]
    if show_curriculum_table:
        overview_rows.append("교육시간 : 하기 표 참고  ※ 강의 시작 10분 전까지 입실 부탁드립니다.")
    else:
        overview_rows.append("교육시간 : 과정별 안내에 따라 별도 확인 부탁드립니다.")
    info_html = build_bullet_rows(info_text.split("\n"), text_color=section_text_color)
    contacts_html = build_contact_rows(edited_contacts, text_color=section_text_color)
    top_logo_html = build_logo_row(logo_image_data_url, logo_position, "top", logo_max_height)
    bottom_logo_html = build_logo_row(logo_image_data_url, logo_position, "bottom", logo_max_height)
    title_margin_top = "14px" if top_logo_html else "4px"

    template_replacements = {
        "{교육명}": f"<span style='font-weight: 800;'>{esc(full_course_name)}</span>",
        "{운영방식}": f"<span style='color: #C1121F; font-weight: 800;'>{esc(delivery_type)}</span>",
        "{1일차}": esc(day1_time),
        "{2일차}": esc(day2_time),
        "{회사명}": esc(company_display),
    }
    welcome_body_html = render_template_html(welcome_body_text, template_replacements)
    time_notice_html = render_template_html(time_notice_text, template_replacements)

    clean_zoom_url = str(zoom_url or "").strip()
    if clean_zoom_url and not re.match(r"^https?://", clean_zoom_url, flags=re.IGNORECASE):
        clean_zoom_url = "https://" + clean_zoom_url
    zoom_button_html = ""
    if clean_zoom_url:
        zoom_button_html = f"""
        <table role="presentation" cellspacing="0" cellpadding="0" border="0" data-zone="overview" style="width: 100%; border-collapse: collapse; margin: 0 0 24px 0;">
            <tr>
                <td style="text-align: center; padding: 0;">
                    <a href="{esc_attr(clean_zoom_url)}" target="_blank" rel="noopener noreferrer" style="display: inline-block; padding: 11px 24px; background-color: {main_color}; color: {curr_header_text_color}; text-decoration: none; border-radius: 8px; font-size: 15px; line-height: 20px; font-weight: 800;">ZOOM 바로가기</a>
                </td>
            </tr>
        </table>
        """

    map_html = ""
    if str(map_image_src or "").strip():
        map_html = f"""
        <table role="presentation" cellspacing="0" cellpadding="0" border="0" style="width: 100%; border-collapse: collapse; margin-top: 12px;">
            <tr>
                <td style="padding: 0;">
                    <img src="{esc_attr(map_image_src)}" alt="교육 장소 지도" style="display: block; width: 100%; max-width: 650px; height: auto; border: 1px solid #D8DEE6;" />
                </td>
            </tr>
        </table>
        """

    return f"""
<div id="mail-content" data-zone="font" style="width: 760px; max-width: 760px; margin: 0 auto; font-family: {font_stack}; background-color: #FFFFFF; color: #222222; border: 1px solid #E5E7EB; box-shadow: 0 10px 28px rgba(15, 23, 42, 0.12);">
    {embedded_font_css}
    <div data-zone="brand" style="height: 18px; background-color: {main_color}; font-size: 0; line-height: 0;">&nbsp;</div>

    <div style="padding: 24px 38px 32px 38px;">
        {top_logo_html}

        <h1 data-zone="title" style="margin: {title_margin_top} 0 20px 0; text-align: center; font-size: 30px; line-height: 38px; font-weight: 800; color: #222222; letter-spacing: -1px;">
            {esc(full_course_name)}
        </h1>

        <div style="border-top: 1px dotted #999999; height: 1px; line-height: 1px; font-size: 0; margin: 0 0 18px 0;">&nbsp;</div>

        <table role="presentation" cellspacing="0" cellpadding="0" border="0" data-zone="welcome time" style="width: 100%; border-collapse: collapse; margin-bottom: 20px;">
            <tr>
                <td style="text-align: center; padding: 0 8px;">
                    <p style="margin: 0 0 12px 0; font-size: 22px; line-height: 30px; font-weight: 800; color: #222222;">
                        {esc(welcome_title)}
                    </p>
                    <p style="margin: 0; font-size: 16px; line-height: 24px; color: #343A40; letter-spacing: -0.3px;">
                        {welcome_body_html}
                    </p>
                    <p style="margin: 12px 0 0 0; font-size: 16px; line-height: 24px; color: #C1121F; font-weight: 800; letter-spacing: -0.3px;">
                        {time_notice_html}
                    </p>
                </td>
            </tr>
        </table>

        <div style="border-top: 1px dotted #999999; height: 1px; line-height: 1px; font-size: 0; margin: 0 0 26px 0;">&nbsp;</div>

        {build_section_title(1, overview_section_title, main_color, curr_header_text_color, "overview")}

        <table role="presentation" cellspacing="0" cellpadding="0" border="0" data-zone="overview" style="width: 100%; border-collapse: collapse; margin: 0 0 18px 0;">
            {build_bullet_rows(overview_rows, text_color=section_text_color)}
        </table>

        {zoom_button_html}

        {curriculum_block_html}

        {build_section_title(2, location_section_title, main_color, curr_header_text_color, "location")}

        <table role="presentation" cellspacing="0" cellpadding="0" border="0" data-zone="location" style="width: 100%; border-collapse: collapse; margin: 0 0 32px 0;">
            {build_bullet_rows(location_lines, text_color=section_text_color)}
            <tr>
                <td colspan="2" style="padding: 0 0 0 26px;">
                    {map_html}
                </td>
            </tr>
        </table>

        {build_section_title(3, notice_section_title, main_color, curr_header_text_color, "notice")}

        <table role="presentation" cellspacing="0" cellpadding="0" border="0" data-zone="notice" style="width: 100%; border-collapse: collapse; margin: 0 0 30px 0;">
            {info_html}
        </table>

        {build_section_title(4, contact_section_title, main_color, curr_header_text_color, "contact")}

        <table role="presentation" cellspacing="0" cellpadding="0" border="0" data-zone="contact" style="width: 100%; border-collapse: collapse; margin: 0 0 10px 0;">
            {contacts_html}
        </table>

        {bottom_logo_html}
    </div>

    <div data-zone="brand" style="height: 18px; background-color: {footer_color}; font-size: 0; line-height: 0;">&nbsp;</div>
</div>
"""




def _ppt_rgb(hex_color: str, fallback: str = "#0088C9"):
    """#RRGGBB 문자열을 python-pptx RGBColor로 변환합니다."""
    try:
        from pptx.dml.color import RGBColor
        clean = str(hex_color or fallback).strip().lstrip("#")
        if len(clean) != 6:
            clean = str(fallback).strip().lstrip("#")
        return RGBColor(int(clean[0:2], 16), int(clean[2:4], 16), int(clean[4:6], 16))
    except Exception:
        from pptx.dml.color import RGBColor
        return RGBColor(0, 136, 201)


def _data_url_to_image_stream(data_url: str):
    """data:image/...;base64 문자열을 BytesIO로 변환합니다."""
    clean = str(data_url or "").strip()
    if not clean.startswith("data:image/") or ";base64," not in clean:
        return None
    try:
        raw = base64.b64decode(clean.split(",", 1)[1])
        return io.BytesIO(raw)
    except Exception:
        return None


def _add_ppt_textbox(slide, text: str, x: float, y: float, w: float, h: float, *,
                     font_size: int = 11, bold: bool = False, color: str = "#222222",
                     align: str = "left", fill: str | None = None, border: str | None = None,
                     radius: bool = False):
    """편집 가능한 PPT 텍스트박스를 추가합니다."""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _ppt_rgb(fill)
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = _ppt_rgb(border)
        shape.line.width = Pt(0.6)
    else:
        shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = str(text or "")
    run.font.size = Pt(font_size)
    run.font.bold = bool(bold)
    run.font.name = "맑은 고딕"
    run.font.color.rgb = _ppt_rgb(color, "#222222")
    return shape


def _add_ppt_multiline_textbox(slide, lines: list[str], x: float, y: float, w: float, h: float, *,
                               font_size: int = 10, color: str = "#343A40", bold_first: bool = False):
    """여러 줄 텍스트를 편집 가능한 PPT 텍스트박스로 추가합니다."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_ANCHOR

    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)

    clean_lines = [str(line or "").strip() for line in lines if str(line or "").strip()]
    if not clean_lines:
        clean_lines = [""]
    for i, line in enumerate(clean_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "맑은 고딕"
        p.font.color.rgb = _ppt_rgb(color)
        p.font.bold = bool(bold_first and i == 0)
        p.space_after = Pt(2)
    return textbox


def _add_ppt_section_title(slide, number: int, title: str, x: float, y: float, w: float,
                           main_color: str, text_color: str) -> float:
    """안내문 섹션 제목을 편집 가능한 도형/텍스트로 추가하고 다음 y좌표를 반환합니다."""
    _add_ppt_textbox(slide, "", x, y, 0.06, 0.30, fill=main_color, border=main_color)
    _add_ppt_textbox(
        slide,
        f"{number}. {title}",
        x + 0.06,
        y,
        min(w - 0.06, 2.7),
        0.30,
        font_size=12,
        bold=True,
        color=text_color,
        fill=main_color,
        border=main_color,
    )
    return y + 0.42


def _add_ppt_bullets(slide, lines: list[str], x: float, y: float, w: float, *,
                     font_size: int = 10, color: str = "#343A40") -> float:
    """하이픈 bullet을 편집 가능한 텍스트로 추가합니다."""
    clean_lines = [normalize_bullet(line) for line in lines if normalize_bullet(line)]
    if not clean_lines:
        clean_lines = ["입력된 내용이 없습니다."]
    line_h = 0.22
    height = max(0.26, len(clean_lines) * line_h + 0.02)
    _add_ppt_multiline_textbox(slide, [f"- {line}" for line in clean_lines], x, y, w, height, font_size=font_size, color=color)
    return y + height + 0.12


def _add_ppt_table(slide, rows: list[dict], columns: list[str], x: float, y: float, w: float,
                   main_color: str, text_color: str, odd_fill: str = "#FFFFFF", even_fill: str = "#FFFFFF") -> float:
    """커리큘럼을 편집 가능한 PPT 표로 추가합니다."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    normalized = normalize_curriculum_for_columns(rows, columns)
    max_rows = max(1, len(normalized))
    row_h = 0.34
    table_h = 0.36 + (max_rows * row_h)
    table_shape = slide.shapes.add_table(max_rows + 1, len(columns), Inches(x), Inches(y), Inches(w), Inches(table_h))
    table = table_shape.table

    # 열 너비 기본값
    weights = []
    for col in columns:
        if any(token in col for token in ["내용", "주제", "과정"]):
            weights.append(2.7)
        elif any(token in col for token in ["비고", "강사"]):
            weights.append(1.15)
        else:
            weights.append(0.85)
    total_weight = sum(weights) or 1
    for idx, weight in enumerate(weights):
        table.columns[idx].width = Inches(w * weight / total_weight)

    for col_idx, col in enumerate(columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _ppt_rgb(main_color)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = "맑은 고딕"
                run.font.size = Pt(9.5)
                run.font.bold = True
                run.font.color.rgb = _ppt_rgb(text_color, "#FFFFFF")

    for row_idx, row in enumerate(normalized, start=1):
        row_fill = odd_fill if row_idx % 2 == 1 else even_fill
        for col_idx, col in enumerate(columns):
            value = _cell_value_by_column(row, col) or "-"
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _ppt_rgb(row_fill, "#FFFFFF")
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_long_text = any(token in col for token in ["내용", "주제", "과정", "비고", "메모"])
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if is_long_text else PP_ALIGN.CENTER
                paragraph.space_after = Pt(0)
                for run in paragraph.runs:
                    run.font.name = "맑은 고딕"
                    run.font.size = Pt(8.5)
                    run.font.bold = bool(is_long_text)
                    run.font.color.rgb = _ppt_rgb("#1F2933")

    return y + table_h + 0.28


def _add_data_url_picture(slide, data_url: str, x: float, y: float, w: float, max_h: float) -> float:
    """data URL 이미지를 PPT에 추가하고 다음 y좌표를 반환합니다."""
    from pptx.util import Inches
    stream = _data_url_to_image_stream(data_url)
    if not stream:
        return y
    try:
        from PIL import Image
        stream.seek(0)
        image = Image.open(stream)
        img_w, img_h = image.size
        stream.seek(0)
        ratio = img_h / img_w if img_w else 0.6
        h = min(max_h, w * ratio)
        slide.shapes.add_picture(stream, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        return y + h + 0.12
    except Exception:
        return y


def _add_logo_to_ppt(slide, logo_data_url: str, logo_position: str, slide_w: float, slide_h: float, max_h: float) -> None:
    """로고 이미지를 사용자가 지정한 위치에 이미지로 추가합니다."""
    from pptx.util import Inches
    stream = _data_url_to_image_stream(logo_data_url)
    if not stream:
        return
    try:
        from PIL import Image
        stream.seek(0)
        image = Image.open(stream)
        img_w, img_h = image.size
        stream.seek(0)
        h = min(max_h, 0.65)
        w = h * (img_w / img_h) if img_h else 1.6
        w = min(w, 2.4)
        x = slide_w - 0.45 - w if "우측" in logo_position else 0.45
        y = slide_h - 0.55 - h if "하단" in logo_position else 0.45
        slide.shapes.add_picture(stream, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    except Exception:
        return


def build_editable_pptx_bytes(
    *,
    company_name: str,
    course_name: str,
    date_range: str,
    day1_time: str,
    day2_time: str,
    place_name: str,
    map_image_src: str,
    delivery_type: str,
    welcome_title: str,
    welcome_body_text: str,
    time_notice_text: str,
    edited_curriculum: list[dict],
    curriculum_columns_text: str,
    curriculum_title: str,
    show_curriculum_table: bool,
    info_text: str,
    edited_contacts: list[dict],
    main_color: str,
    footer_color: str,
    curr_header_text_color: str,
    logo_image_data_url: str,
    logo_position: str,
    logo_max_height: int,
    zoom_url: str = "",
    overview_section_title: str = "교육 개요",
    location_section_title: str = "교육 장소",
    notice_section_title: str = "안내 사항",
    contact_section_title: str = "관련 문의",
    curriculum_odd_row_color: str = "#FFFFFF",
    curriculum_even_row_color: str = "#FFFFFF",
) -> bytes:
    """현재 입력값을 기반으로 텍스트/도형/표가 살아있는 PPTX 1장을 생성합니다.

    로고와 지도는 원본 이미지로 들어가지만, 제목/본문/섹션/시간표/문의처는
    PowerPoint에서 직접 수정 가능한 텍스트·도형·표로 생성됩니다.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE

    company_display = str(company_name or "").strip()
    full_course_name = f"{company_display} {course_name}".strip()
    curriculum_columns = parse_curriculum_columns(curriculum_columns_text)
    curriculum_rows = normalize_curriculum_for_columns(edited_curriculum, curriculum_columns)
    info_lines = [normalize_bullet(line) for line in str(info_text or "").split("\n") if normalize_bullet(line)]
    contact_lines = []
    for contact in to_records(edited_contacts):
        role = str(contact.get("role", "") or "").strip()
        phone = str(contact.get("phone", "") or "").strip()
        if role or phone:
            contact_lines.append(f"{role} ({phone})" if phone else role)

    map_h_est = 2.15 if str(map_image_src or "").strip() else 0
    curr_h_est = (0.8 + len(curriculum_rows) * 0.34) if show_curriculum_table else 0
    info_h_est = max(0.45, len(info_lines) * 0.22 + 0.1)
    contact_h_est = max(0.45, len(contact_lines) * 0.22 + 0.1)
    slide_w = 7.6
    slide_h = max(10.8, 5.25 + curr_h_est + map_h_est + info_h_est + contact_h_est)
    slide_h = min(slide_h, 20.5)

    prs = Presentation()
    prs.slide_width = Inches(slide_w)
    prs.slide_height = Inches(slide_h)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _ppt_rgb("#FFFFFF")

    # 외곽 배경/상하단 브랜드 바
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(slide_w), Inches(slide_h)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = _ppt_rgb("#FFFFFF")
    slide.shapes[-1].line.color.rgb = _ppt_rgb("#E5E7EB")
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(slide_w), Inches(0.18)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = _ppt_rgb(main_color)
    slide.shapes[-1].line.color.rgb = _ppt_rgb(main_color)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(slide_h - 0.18), Inches(slide_w), Inches(0.18)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = _ppt_rgb(footer_color)
    slide.shapes[-1].line.color.rgb = _ppt_rgb(footer_color)

    _add_logo_to_ppt(slide, logo_image_data_url, logo_position, slide_w, slide_h, max(0.28, float(logo_max_height or 52) / 80))

    x = 0.38
    w = slide_w - 0.76
    y = 0.58
    _add_ppt_textbox(slide, full_course_name, x, y, w, 0.42, font_size=19, bold=True, color="#222222", align="center")
    y += 0.62
    _add_ppt_textbox(slide, welcome_title, x, y, w, 0.32, font_size=14, bold=True, color="#222222", align="center")
    y += 0.34

    replacements_plain = {
        "{교육명}": full_course_name,
        "{운영방식}": delivery_type,
        "{1일차}": day1_time,
        "{2일차}": day2_time,
        "{회사명}": company_display,
    }
    welcome_plain = str(welcome_body_text or "")
    time_plain = str(time_notice_text or "")
    for token, value in replacements_plain.items():
        welcome_plain = welcome_plain.replace(token, str(value or ""))
        time_plain = time_plain.replace(token, str(value or ""))
    _add_ppt_multiline_textbox(slide, welcome_plain.split("\n"), x + 0.25, y, w - 0.5, 0.55, font_size=10, color="#343A40")
    y += 0.58
    _add_ppt_textbox(slide, time_plain, x + 0.25, y, w - 0.5, 0.30, font_size=10, bold=True, color="#C1121F", align="center")
    y += 0.52

    # 1. 교육 개요
    y = _add_ppt_section_title(slide, 1, overview_section_title, x, y, w, main_color, curr_header_text_color)
    overview_lines = [
        f"교육명 : {full_course_name}",
        f"교육일정 : {date_range}  *집합 교육 기준",
        "교육시간 : 하기 표 참고  ※ 강의 시작 10분 전까지 입실 부탁드립니다." if show_curriculum_table else "교육시간 : 과정별 안내에 따라 별도 확인 부탁드립니다.",
    ]
    y = _add_ppt_bullets(slide, overview_lines, x + 0.16, y, w - 0.25, font_size=9.5)

    clean_zoom_url = str(zoom_url or "").strip()
    if clean_zoom_url:
        if not re.match(r"^https?://", clean_zoom_url, flags=re.IGNORECASE):
            clean_zoom_url = "https://" + clean_zoom_url
        btn = _add_ppt_textbox(slide, "ZOOM 바로가기", x + 0.16, y, 1.65, 0.34, font_size=10, bold=True, color=curr_header_text_color, fill=main_color, border=main_color, align="center", radius=True)
        try:
            btn.click_action.hyperlink.address = clean_zoom_url
        except Exception:
            pass
        y += 0.52

    if show_curriculum_table:
        if str(curriculum_title or "").strip():
            _add_ppt_textbox(slide, curriculum_title, x, y, w, 0.25, font_size=10, bold=True, color="#222222")
            y += 0.30
        y = _add_ppt_table(slide, curriculum_rows, curriculum_columns, x, y, w, main_color, curr_header_text_color, curriculum_odd_row_color, curriculum_even_row_color)

    # 2. 교육 장소
    y = _add_ppt_section_title(slide, 2, location_section_title, x, y, w, main_color, curr_header_text_color)
    location_lines = [line.strip() for line in str(place_name or "").splitlines() if line.strip()]
    y = _add_ppt_bullets(slide, location_lines, x + 0.16, y, w - 0.25, font_size=9.5)
    if str(map_image_src or "").strip():
        y = _add_data_url_picture(slide, map_image_src, x + 0.25, y, min(6.5, w - 0.5), 2.5)

    # 3. 안내 사항
    y = _add_ppt_section_title(slide, 3, notice_section_title, x, y, w, main_color, curr_header_text_color)
    y = _add_ppt_bullets(slide, info_lines, x + 0.16, y, w - 0.25, font_size=9.5)

    # 4. 관련 문의
    y = _add_ppt_section_title(slide, 4, contact_section_title, x, y, w, main_color, curr_header_text_color)
    _add_ppt_bullets(slide, contact_lines, x + 0.16, y, w - 0.25, font_size=9.5)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

def build_preview_component_html(final_mail_html: str, preview_scale: float, font_stack: str, export_base_name: str = "education_notice_photo_style", app_theme: str = "다크 모드") -> str:
    scale = max(0.45, min(1.0, float(preview_scale)))
    export_base_name_js = json.dumps(str(export_base_name or "education_notice_photo_style"), ensure_ascii=False)
    is_light_theme = str(app_theme or "") == "화이트 모드"
    tool_btn_bg = "#7b7b7b" if is_light_theme else "#111111"
    tool_btn_border = "#c9c9c9" if is_light_theme else "#111111"
    tool_btn_hover = "#7b7b7b" if is_light_theme else "#333333"
    return (
        f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8" />
    <script src="https://cdnjs.cloudflare.com/ajax/libs/html2canvas/1.4.1/html2canvas.min.js"></script>
    <style>
        body {{
            margin: 0;
            padding: 14px 0 34px 0;
            background: #EEF2F6;
            font-family: {font_stack};
        }}
        .toolbar {{
            width: min(760px, calc(100vw - 28px));
            margin: 0 auto 10px auto;
            display: grid;
            grid-template-columns: repeat(4, 1fr);
            gap: 8px;
            position: sticky;
            top: 0;
            z-index: 5;
            background: #EEF2F6;
            padding-bottom: 8px;
        }}
        .tool-btn {{
            border: 1px solid {tool_btn_border};
            border-radius: 12px;
            min-height: 40px;
            padding: 5px 6px;
            font-size: 12px;
            line-height: 15px;
            font-weight: 800;
            color: #FFFFFF;
            background: {tool_btn_bg};
            cursor: pointer;
        }}
        .tool-btn:hover {{ background: {tool_btn_hover}; }}
        .note {{
            width: min(760px, calc(100vw - 28px));
            margin: 0 auto 10px auto;
            font-size: 12px;
            line-height: 18px;
            color: #737373;
        }}
        .preview-area {{
            zoom: {scale};
        }}
        @supports not (zoom: 1) {{
            .preview-area {{
                transform: scale({scale});
                transform-origin: top center;
            }}
        }}
    </style>
</head>
<body>
    <div class="toolbar">
        <button class="tool-btn" onclick="copyMailContent()">📋 아웃룩<br>서식 복사</button>
        <button class="tool-btn" onclick="copyHtmlCode()">🖱️ HTML<br>코드 복사</button>
        <button class="tool-btn" onclick="downloadImage('png')">🖼️ PNG<br>저장</button>
        <button class="tool-btn" onclick="downloadImage('jpeg')">🖼️ JPG<br>저장</button>
    </div>
    <div class="note">PPT 다운로드는 아래 “편집 가능한 PPT 다운로드” 버튼을 사용하세요.</div>
    <div class="preview-area">
"""
        + final_mail_html
        + """
    </div>
    <script>
        const EXPORT_BASE_NAME = """ + export_base_name_js + """;

        function getContent() {
            return document.getElementById('mail-content');
        }

        async function copyMailContent() {
            const content = getContent();
            if (!content) {
                alert('복사할 대상을 찾지 못했습니다.');
                return;
            }

            const htmlContent = content.outerHTML;
            const textContent = content.innerText;

            try {
                if (navigator.clipboard && window.ClipboardItem) {
                    await navigator.clipboard.write([
                        new ClipboardItem({
                            'text/html': new Blob([htmlContent], { type: 'text/html' }),
                            'text/plain': new Blob([textContent], { type: 'text/plain' })
                        })
                    ]);
                    alert('사진형 안내문 서식이 복사되었습니다. 아웃룩 본문에 Ctrl+V 하세요.');
                    return;
                }
            } catch (error) {
                console.log(error);
            }

            try {
                const range = document.createRange();
                range.selectNodeContents(content);
                const selection = window.getSelection();
                selection.removeAllRanges();
                selection.addRange(range);
                document.execCommand('copy');
                selection.removeAllRanges();
                alert('사진형 안내문 서식이 복사되었습니다. 아웃룩 본문에 Ctrl+V 하세요.');
            } catch (error) {
                alert('복사에 실패했습니다. 미리보기 영역을 직접 드래그해서 복사해 주세요.');
            }
        }

        async function copyHtmlCode() {
            const content = getContent();
            if (!content) {
                alert('복사할 HTML 코드를 찾지 못했습니다.');
                return;
            }

            const htmlCode = content.outerHTML;

            try {
                if (navigator.clipboard && navigator.clipboard.writeText) {
                    await navigator.clipboard.writeText(htmlCode);
                    alert('HTML 코드가 클립보드에 복사되었습니다.');
                    return;
                }
            } catch (error) {
                console.log(error);
            }

            try {
                const textarea = document.createElement('textarea');
                textarea.value = htmlCode;
                textarea.setAttribute('readonly', '');
                textarea.style.position = 'fixed';
                textarea.style.left = '-9999px';
                document.body.appendChild(textarea);
                textarea.select();
                document.execCommand('copy');
                document.body.removeChild(textarea);
                alert('HTML 코드가 클립보드에 복사되었습니다.');
            } catch (error) {
                alert('HTML 코드 복사에 실패했습니다. Raw HTML 탭에서 직접 복사해 주세요.');
            }
        }

        async function renderNoticeCanvas() {
            const source = getContent();
            if (!source) {
                throw new Error('저장할 안내문을 찾지 못했습니다.');
            }
            if (typeof html2canvas === 'undefined') {
                throw new Error('이미지 저장 모듈을 불러오지 못했습니다. 인터넷 연결 후 다시 시도해 주세요.');
            }

            const clone = source.cloneNode(true);
            clone.id = 'mail-content-export';
            clone.style.margin = '0';

            const offscreen = document.createElement('div');
            offscreen.style.position = 'fixed';
            offscreen.style.left = '-10000px';
            offscreen.style.top = '0';
            offscreen.style.width = source.offsetWidth + 'px';
            offscreen.style.background = '#ffffff';
            offscreen.style.zIndex = '-1';
            offscreen.appendChild(clone);
            document.body.appendChild(offscreen);

            try {
                return await html2canvas(clone, {
                    scale: 2,
                    backgroundColor: '#ffffff',
                    useCORS: true,
                    allowTaint: false,
                    logging: false,
                    scrollX: 0,
                    scrollY: 0,
                    windowWidth: 1200,
                    windowHeight: clone.scrollHeight + 100
                });
            } finally {
                document.body.removeChild(offscreen);
            }
        }

        async function downloadImage(format) {
            try {
                const canvas = await renderNoticeCanvas();
                const mimeType = format === 'jpeg' ? 'image/jpeg' : 'image/png';
                const extension = format === 'jpeg' ? 'jpg' : 'png';
                const dataUrl = canvas.toDataURL(mimeType, 0.95);
                const link = document.createElement('a');
                link.download = EXPORT_BASE_NAME + '.' + extension;
                link.href = dataUrl;
                document.body.appendChild(link);
                link.click();
                document.body.removeChild(link);
            } catch (error) {
                console.error(error);
                alert(error.message || '이미지 저장에 실패했습니다.');
            }
        }

    </script>
</body>
</html>
"""
    )


# -----------------------------
# 쿼리 파라미터 기반 스포이드 색상 반영
# -----------------------------
init_defaults()
migrate_curriculum_day_defaults()
apply_pending_location_update()
poll_split_naver_jobs()

picked_color = get_query_param("picked_color")
picked_target = get_query_param("picked_target")
picked_token = get_query_param("picked_token")
if is_valid_hex_color(picked_color) and picked_token and st.session_state.get("last_picked_token") != picked_token:
    if picked_target == "main":
        st.session_state.main_color_picker = picked_color
    elif picked_target == "footer":
        st.session_state.footer_color_picker = picked_color
        st.session_state.use_custom_footer_color = True
    st.session_state.last_picked_token = picked_token
    try:
        st.query_params.clear()
    except Exception:
        try:
            st.experimental_set_query_params()
        except Exception:
            pass
    st.rerun()


SIDE_PANEL_HEIGHT = 920
PREVIEW_IFRAME_HEIGHT = 690

# -----------------------------
# 입력 / 출력 레이아웃
# -----------------------------
col_input, col_output = st.columns([0.9, 1.25], gap="large")

with col_input:
    st.subheader("입력 설정")

    with st.container(height=SIDE_PANEL_HEIGHT, border=False):

        # 1. 기본 정보
        with st.container(border=True):
            st.markdown("#### 기본 정보")
            company_name = st.text_input(
                "회사 이름",
                key="company_name",
                help="입력한 회사 이름이 그대로 표시됩니다. '그룹'은 자동으로 붙지 않습니다.",
            )
            course_name = st.text_input("교육 과정명", key="course_name")
            delivery_mode = st.selectbox(
                "운영 방식",
                ["대면", "비대면", "직접입력"],
                key="delivery_mode",
            )
            if delivery_mode == "직접입력":
                delivery_type = st.text_input(
                    "운영 방식 직접 입력",
                    key="delivery_custom",
                    placeholder="예: 하이브리드, 온라인 실시간, 집합",
                ).strip()
            else:
                delivery_type = delivery_mode

            delivery_type_lower = str(delivery_type or "").strip().lower()
            is_online_delivery = delivery_mode == "비대면" or any(
                token in delivery_type_lower for token in ["비대면", "온라인", "zoom", "줌", "하이브리드"]
            )
            if is_online_delivery:
                zoom_url = st.text_input(
                    "ZOOM/온라인 과정 바로가기 링크",
                    key="zoom_url",
                    placeholder="https://zoom.us/j/...",
                    help="비대면 또는 하이브리드 과정일 때 안내문에 바로가기 버튼을 표시합니다.",
                )
            else:
                zoom_url = ""

            with st.expander("섹션 제목 수정", expanded=False):
                section_title_col1, section_title_col2 = st.columns(2)
                with section_title_col1:
                    overview_section_title = st.text_input("1번 제목", key="overview_section_title")
                    notice_section_title = st.text_input("3번 제목", key="notice_section_title")
                with section_title_col2:
                    location_section_title = st.text_input("2번 제목", key="location_section_title")
                    contact_section_title = st.text_input("4번 제목", key="contact_section_title")

        # 2. 상단 안내 문구
        with st.container(border=True):
            st.markdown("#### 상단 안내 문구")
            st.caption("사용 가능 치환값: {교육명}, {운영방식}, {1일차}, {2일차}, {회사명}")
            welcome_title = st.text_input("환영 제목", key="welcome_title")
            welcome_body_text = st.text_area("상단 안내 본문", key="welcome_body_text", height=110)
            time_notice_text = st.text_area("교육 시작시간 강조 문구", key="time_notice_text", height=80)
            date_range = st.text_input("교육 일정 (예: 6/16(월)~6/17(화))", key="date_range")
            col_t1, col_t2 = st.columns(2)
            with col_t1:
                day1_time = st.text_input("1일차 시작 시간", key="day1_time")
            with col_t2:
                day2_time = st.text_input("2일차 시작 시간", key="day2_time")

        # 3. 브랜드 컬러 + 로고/폰트
        with st.container(border=True):
            st.markdown("#### 브랜드 컬러")
            col_c1, col_c2, col_c3 = st.columns(3)
            with col_c1:
                main_color = st.color_picker("메인 컬러", key="main_color_picker")
                eyedropper_component("main", "컬러 스포이드")
            with col_c2:
                use_custom_footer_color = st.checkbox("하단 컬러 별도 지정", key="use_custom_footer_color")
                if use_custom_footer_color:
                    footer_color = st.color_picker("하단 컬러", key="footer_color_picker")
                    eyedropper_component("footer", "컬러 스포이드")
                else:
                    footer_color = main_color
                    st.caption("하단 바도 메인 컬러를 사용합니다.")
            with col_c3:
                curr_header_text_color = st.color_picker(
                    "컬러박스 내 글자색",
                    key="curr_header_text_color_picker",
                )

            section_text_color = "#000000"

            st.markdown("---")
            st.markdown("#### 로고 / 폰트")
            logo_file = st.file_uploader(
                "회사 로고 이미지 첨부",
                type=["png", "jpg", "jpeg", "gif", "webp"],
                help="첨부하지 않으면 로고 영역은 비워집니다. 회사명 텍스트 로고는 자동으로 표시하지 않습니다.",
                key="logo_file_uploader",
            )
            logo_image_data_url = file_to_data_url(logo_file)
            col_l1, col_l2 = st.columns(2)
            with col_l1:
                logo_position = st.selectbox(
                    "로고 위치",
                    ["우측 상단", "좌측 상단", "우측 하단", "좌측 하단"],
                    key="logo_position",
                    help="로고 이미지를 첨부한 경우에만 적용됩니다.",
                )
            with col_l2:
                logo_max_height = st.slider("로고 크기 조절", min_value=28, max_value=90, step=2, key="logo_max_height")

            embedded_font_css = ""
            fonts_dir = get_fonts_dir()
            try:
                fonts_dir.mkdir(parents=True, exist_ok=True)
            except OSError:
                pass
            font_files = list_font_files(str(fonts_dir))
            if font_files:
                selected_font_file = st.selectbox(
                    "폰트 선택",
                    font_files,
                    format_func=lambda x: Path(x).name,
                    key="font_folder_select",
                )
                primary_font = font_family_from_file(selected_font_file)
                embedded_font_css = build_embedded_font_css(selected_font_file, primary_font)
            else:
                primary_font = "Malgun Gothic"
                st.caption("assets/fonts 폴더에 폰트 파일이 없어서 기본 폰트를 사용합니다.")
            font_stack = css_font_stack(primary_font)

        # 4. 커리큘럼
        with st.container(border=True):
            st.markdown("#### 커리큘럼")
            show_curriculum_table = st.checkbox(
                "교육 시간표 표시",
                key="show_curriculum_table",
                help="체크를 해제하면 최종 안내문에서 커리큘럼 표를 표시하지 않습니다.",
            )
            curriculum_title = st.text_input("커리큘럼 표 이름", key="curriculum_title")
            if not show_curriculum_table:
                st.caption("교육 시간표 표시를 해제했습니다. 입력값은 유지되지만 최종 안내문에는 표가 나오지 않습니다.")

            with st.expander("엑셀 시간표 불러오기", expanded=False):
                st.caption("첫 번째 시트에서 Day·시간·교육 내용·강사/비고 열을 찾아 불러옵니다. 지원 형식: .xlsx, .xlsm")
                curriculum_excel_file = st.file_uploader(
                    "시간표 엑셀 파일",
                    type=["xlsx", "xlsm"],
                    key="curriculum_excel_uploader",
                )
                if st.button("엑셀 시간표 적용", key="apply_curriculum_excel", use_container_width=True):
                    try:
                        current_columns = [
                            str(row.get("column_name", "") or "").strip()
                            for row in to_records(st.session_state.get("curriculum_column_defs", []))
                            if str(row.get("column_name", "") or "").strip()
                        ]
                        imported_rows, imported_columns, import_message = parse_curriculum_excel(
                            curriculum_excel_file,
                            current_columns,
                        )
                        st.session_state.curriculum = imported_rows
                        st.session_state.curriculum_column_defs = [
                            {"column_name": column} for column in imported_columns
                        ]
                        st.session_state.curriculum_columns_text = ", ".join(imported_columns)
                        st.session_state.excel_import_message = import_message
                        st.session_state.pop("curr_editor", None)
                        st.session_state.pop("curriculum_column_defs_editor", None)
                        for state_key in list(st.session_state.keys()):
                            if state_key.startswith("white_grid_curr_editor_") or state_key.startswith("white_grid_curriculum_column_defs_editor_"):
                                st.session_state.pop(state_key, None)
                        st.rerun()
                    except Exception as exc:
                        st.session_state.excel_import_message = str(exc)
                if st.session_state.get("excel_import_message"):
                    st.caption(st.session_state.excel_import_message)

            with st.expander("시간표 음영 색상", expanded=False):
                shade_col1, shade_col2 = st.columns(2)
                with shade_col1:
                    curriculum_odd_row_color = st.color_picker(
                        "홀수 행 음영",
                        key="curriculum_odd_row_color_picker",
                    )
                with shade_col2:
                    curriculum_even_row_color = st.color_picker(
                        "짝수 행 음영",
                        key="curriculum_even_row_color_picker",
                    )

            with st.expander("표 헤더 편집", expanded=False):
                if st.session_state.get("app_theme") == "화이트 모드":
                    header_editor_data = render_white_theme_grid_editor(
                        st.session_state.curriculum_column_defs,
                        ["column_name"],
                        key_prefix="curriculum_column_defs_editor",
                        column_labels={"column_name": "표 헤더명"},
                        min_rows=1,
                        add_button_label="표 헤더 추가",
                    )
                else:
                    header_editor_data = st.data_editor(
                        st.session_state.curriculum_column_defs,
                        num_rows="dynamic",
                        hide_index=True,
                        column_order=["column_name"],
                        column_config={"column_name": st.column_config.TextColumn("표 헤더명", required=True)},
                        width="stretch",
                        key="curriculum_column_defs_editor",
                    )
                header_records = to_records(header_editor_data)
                curriculum_columns = [str(row.get("column_name", "") or "").strip() for row in header_records]
                curriculum_columns = [col for col in curriculum_columns if col]
                if not curriculum_columns:
                    curriculum_columns = ["시간", "Day", "교육 내용", "강사/비고"]
                st.session_state.curriculum_column_defs = [{"column_name": col} for col in curriculum_columns]
                st.session_state.curriculum_columns_text = ", ".join(curriculum_columns)

            curriculum_columns = [str(row.get("column_name", "") or "").strip() for row in st.session_state.curriculum_column_defs]
            curriculum_columns = [col for col in curriculum_columns if col]
            if not curriculum_columns:
                curriculum_columns = parse_curriculum_columns(st.session_state.get("curriculum_columns_text", "Day, 시간, 교육 내용, 강사/비고"))
                st.session_state.curriculum_column_defs = [{"column_name": col} for col in curriculum_columns]
            curriculum_columns_text = ", ".join(curriculum_columns)

            curriculum_editor_data = normalize_curriculum_for_columns(st.session_state.curriculum, curriculum_columns)
            if st.session_state.get("app_theme") == "화이트 모드":
                edited_curr = render_white_theme_grid_editor(
                    curriculum_editor_data,
                    curriculum_columns,
                    key_prefix="curr_editor",
                    column_labels={column: column for column in curriculum_columns},
                    min_rows=1,
                    add_button_label="커리큘럼 행 추가",
                )
            else:
                edited_curr = st.data_editor(
                    curriculum_editor_data,
                    num_rows="dynamic",
                    column_order=curriculum_columns,
                    column_config={column: column for column in curriculum_columns},
                    width="stretch",
                    key="curr_editor",
                )
            st.session_state.curriculum = to_records(edited_curr)

        # 5. 장소
        with st.container(border=True):
            st.markdown("#### 장소")

            place_name = st.text_input("교육 장소명 / 검색어", key="place_name")
            naver_map_query = quote(str(place_name or "").strip())
            naver_map_url = f"https://map.naver.com/p/search/{naver_map_query}?c=15.00,0,0,0,dh" if naver_map_query else "https://map.naver.com/p/search/"

            st.button(
                "📍 주소 자동 입력",
                use_container_width=True,
                on_click=apply_first_naver_local_result,
                help="장소 후보를 검색하고 첫 번째 후보 주소를 바로 입력합니다.",
            )

            local_results = st.session_state.get("naver_local_results", []) or []
            if local_results:
                if st.session_state.get("selected_naver_local_index", 0) not in list(range(len(local_results))):
                    st.session_state.selected_naver_local_index = 0

                def _local_result_label(index: int) -> str:
                    item = local_results[index]
                    title = strip_html_tags(item.get("title", ""))
                    road = strip_html_tags(item.get("roadAddress", ""))
                    addr = strip_html_tags(item.get("address", ""))
                    return f"{title} | {road or addr}"

                col_candidate1, col_candidate2 = st.columns([1.45, 0.55])
                with col_candidate1:
                    selected_local_index = st.selectbox(
                        "검색 후보",
                        options=list(range(len(local_results))),
                        format_func=_local_result_label,
                        key="selected_naver_local_index",
                    )
                with col_candidate2:
                    st.markdown("<div style='height: 28px;'></div>", unsafe_allow_html=True)
                    st.button(
                        "후보 사용",
                        use_container_width=True,
                        on_click=apply_naver_local_item,
                        args=(local_results[selected_local_index],),
                    )

            last_address_fetch_message = st.session_state.get("last_address_fetch_message", "")
            if last_address_fetch_message:
                st.caption(last_address_fetch_message)

            display_location_text = st.text_area(
                "안내문 표시 내용",
                key="display_location_text",
                height=68,
                help="최종 안내문 교육 장소 영역에 들어갈 문구입니다. 예: 멀티캠퍼스 선릉  (서울 강남구 선릉로 428)",
            )

            st.markdown("##### 지도 이미지")
            uploaded_map_data_url = ""
            pasted_map_data_url = normalize_image_data_url(st.session_state.get("pasted_map_data_url", ""))
            captured_map_path = st.session_state.get("captured_map_file_path", "") or st.session_state.get("last_captured_map_file", "")
            captured_map_data_url = file_path_to_data_url(captured_map_path) or st.session_state.get("captured_map_data_url", "")

            st.markdown(
                """
                <style>
                textarea[aria-label="클립보드 이미지 데이터"] {
                    height: 1px !important;
                    min-height: 1px !important;
                    padding: 0 !important;
                    border: 0 !important;
                    opacity: 0.01 !important;
                }
                div.st-key-paste_refresh_trigger {
                    height: 1px !important;
                    min-height: 1px !important;
                    overflow: hidden !important;
                    margin: 0 !important;
                    padding: 0 !important;
                }
                div.st-key-paste_refresh_trigger button {
                    height: 1px !important;
                    min-height: 1px !important;
                    padding: 0 !important;
                    opacity: 0.01 !important;
                }
                </style>
                """,
                unsafe_allow_html=True,
            )

            col_paste, col_map_tools = st.columns([1.35, 0.65])
            with col_paste:
                clipboard_map_paste_component()
            with col_map_tools:
                st.markdown(
                    f'<a class="naver-map-button" href="{naver_map_url}" target="_blank" rel="noopener noreferrer">🗺️ 네이버 지도</a>',
                    unsafe_allow_html=True,
                )

            map_file = st.file_uploader(
                "지도 이미지 파일 첨부",
                type=["png", "jpg", "jpeg", "gif", "webp"],
                help="이미지 파일을 저장해둔 경우에는 여기로 첨부해 주세요. 첨부 파일이 붙여넣기 이미지보다 우선 적용됩니다.",
                key="map_file_uploader",
            )
            uploaded_map_data_url = file_to_data_url(map_file)

            pasted_map_input = st.text_area(
                "클립보드 이미지 데이터",
                key="pasted_map_data_url",
                height=1,
                placeholder="클립보드 이미지 데이터",
                label_visibility="collapsed",
            )
            if st.button("붙여넣기 반영", key="paste_refresh_trigger"):
                st.session_state.last_map_capture_message = "붙여넣은 지도 이미지를 적용했습니다."
            pasted_map_data_url = normalize_image_data_url(pasted_map_input)

            with st.expander("캡처가 안 될 때(저속 보조 기능)", expanded=False):
                col_map_cap1, col_map_cap2 = st.columns([1, 1])
                with col_map_cap1:
                    if st.button("📸 자동 캡처 실행", use_container_width=True):
                        with st.spinner("네이버 지도 검색 화면을 열고 캡처 중입니다. 환경에 따라 20~40초 정도 걸릴 수 있습니다..."):
                            ok, message = capture_naver_map_region(place_name)
                        st.session_state.last_map_capture_message = message
                        if ok:
                            st.rerun()
                        else:
                            st.error(message)
                with col_map_cap2:
                    if st.button("🧹 자동 캡처 제거", use_container_width=True):
                        st.session_state.captured_map_data_url = ""
                        st.session_state.captured_map_file_path = ""
                        st.session_state.pop("last_captured_map_file", None)
                        st.session_state.pop("captured_map_files", None)
                        st.session_state.last_map_capture_message = "자동 캡처 지도를 제거했습니다."
                        st.rerun()

            captured_map_path = st.session_state.get("captured_map_file_path", "") or st.session_state.get("last_captured_map_file", "")
            captured_map_data_url = file_path_to_data_url(captured_map_path) or st.session_state.get("captured_map_data_url", "")
            map_image_src = uploaded_map_data_url or pasted_map_data_url or captured_map_data_url

            last_map_capture_message = st.session_state.get("last_map_capture_message", "")
            if last_map_capture_message:
                st.caption(last_map_capture_message)

            if uploaded_map_data_url:
                st.image(uploaded_map_data_url, caption="첨부한 지도 이미지", use_container_width=True)
            elif pasted_map_data_url:
                st.image(pasted_map_data_url, caption="붙여넣은 지도 이미지", use_container_width=True)
            elif captured_map_data_url:
                st.image(captured_map_data_url, caption="자동 캡처된 지도 이미지", use_container_width=True)


            if map_image_src:
                st.caption("지도 이미지가 적용되어 있습니다.")
            else:
                st.caption("지도 이미지를 붙여넣거나 첨부하면 안내문에 반영됩니다.")

        # 6. 안내 사항 / 문의처
        with st.container(border=True):
            st.markdown("#### 안내 사항 / 문의처")
            info_text = st.text_area("안내 사항 문구", key="info_text", height=120)

            if st.session_state.get("app_theme") == "화이트 모드":
                edited_contacts = render_white_theme_grid_editor(
                    st.session_state.contacts,
                    ["role", "phone"],
                    key_prefix="contacts_editor",
                    column_labels={"role": "소속 + 이름 + 직급", "phone": "연락처"},
                    min_rows=1,
                    add_button_label="문의처 행 추가",
                )
            else:
                edited_contacts = st.data_editor(
                    st.session_state.contacts,
                    num_rows="dynamic",
                    column_config={"role": "소속 + 이름 + 직급", "phone": "연락처"},
                    width="stretch",
                    key="contacts_editor",
                )



final_mail_html = build_final_mail_html(
    company_name=company_name,
    course_name=course_name,
    date_range=date_range,
    day1_time=day1_time,
    day2_time=day2_time,
    place_name=display_location_text,
    road_address="",
    map_image_src=map_image_src,
    delivery_type=delivery_type,
    welcome_title=welcome_title,
    welcome_body_text=welcome_body_text,
    time_notice_text=time_notice_text,
    edited_curriculum=edited_curr,
    curriculum_columns_text=curriculum_columns_text,
    curriculum_title=curriculum_title,
    show_curriculum_table=show_curriculum_table,
    info_text=info_text,
    edited_contacts=edited_contacts,
    main_color=main_color,
    footer_color=footer_color,
    section_text_color=section_text_color,
    curr_header_text_color=curr_header_text_color,
    logo_image_data_url=logo_image_data_url,
    logo_position=logo_position,
    logo_max_height=logo_max_height,
    font_stack=font_stack,
    embedded_font_css=embedded_font_css,
    zoom_url=zoom_url,
    overview_section_title=overview_section_title,
    location_section_title=location_section_title,
    notice_section_title=notice_section_title,
    contact_section_title=contact_section_title,
    curriculum_odd_row_color=curriculum_odd_row_color,
    curriculum_even_row_color=curriculum_even_row_color,
)


export_file_base = sanitize_file_name(st.session_state.get("export_file_name", DEFAULT_VALUES["export_file_name"]))
download_html_document = build_download_html_document(
    final_mail_html=final_mail_html,
    font_stack=font_stack,
    page_title=export_file_base,
)
# Windows 기본 브라우저/일부 편집기에서 한글이 깨지는 것을 줄이기 위해 UTF-8 BOM을 포함합니다.
download_html_bytes = ("\ufeff" + download_html_document).encode("utf-8")


try:
    editable_pptx_bytes = build_editable_pptx_bytes(
        company_name=company_name,
        course_name=course_name,
        date_range=date_range,
        day1_time=day1_time,
        day2_time=day2_time,
        place_name=display_location_text,
        map_image_src=map_image_src,
        delivery_type=delivery_type,
        welcome_title=welcome_title,
        welcome_body_text=welcome_body_text,
        time_notice_text=time_notice_text,
        edited_curriculum=to_records(edited_curr),
        curriculum_columns_text=curriculum_columns_text,
        curriculum_title=curriculum_title,
        show_curriculum_table=show_curriculum_table,
        info_text=info_text,
        edited_contacts=to_records(edited_contacts),
        main_color=main_color,
        footer_color=footer_color,
        curr_header_text_color=curr_header_text_color,
        logo_image_data_url=logo_image_data_url,
        logo_position=logo_position,
        logo_max_height=logo_max_height,
        zoom_url=zoom_url,
        overview_section_title=overview_section_title,
        location_section_title=location_section_title,
        notice_section_title=notice_section_title,
        contact_section_title=contact_section_title,
        curriculum_odd_row_color=curriculum_odd_row_color,
        curriculum_even_row_color=curriculum_even_row_color,
    )
    editable_pptx_error = ""
except Exception as exc:
    editable_pptx_bytes = b""
    editable_pptx_error = str(exc)


with col_output:
    st.subheader("출력 미리보기")

    with st.container(height=SIDE_PANEL_HEIGHT, border=False):
        tab_preview, tab_code = st.tabs(["👀 압축 미리보기", "💻 Raw HTML 소스코드"])

        with tab_preview:
            preview_scale_percent_for_render = st.session_state.get(
                "preview_scale_percent",
                DEFAULT_VALUES["preview_scale_percent"],
            )
            preview_scale = preview_scale_percent_for_render / 100

            components.html(
                build_preview_component_html(
                    final_mail_html,
                    preview_scale=preview_scale,
                    font_stack=font_stack,
                    export_base_name=export_file_base,
                    app_theme=st.session_state.get("app_theme", DEFAULT_VALUES["app_theme"]),
                ),
                height=PREVIEW_IFRAME_HEIGHT,
                scrolling=True,
            )

            with st.container(border=True):
                st.markdown("#### 내보내기 / 미리보기 설정")
                st.text_input(
                    "내보낼 파일 이름",
                    key="export_file_name",
                    help="확장자는 자동으로 붙습니다. 예: 멀티캠퍼스_입문교육_안내문",
                )
                col_p1, col_p2 = st.columns([1.4, 1])
                with col_p1:
                    st.caption("미리보기는 축소해서 보여주고, 저장 이미지는 원본 크기로 생성됩니다.")
                with col_p2:
                    st.slider("미리보기 배율", min_value=50, max_value=100, step=2, key="preview_scale_percent")

                col_d1, col_d2, col_d3 = st.columns([1, 1, 1])
                with col_d1:
                    st.download_button(
                        label="⬇️ HTML 파일 다운로드",
                        data=download_html_bytes,
                        file_name=f"{export_file_base}.html",
                        mime="text/html; charset=utf-8",
                        use_container_width=True,
                    )
                with col_d2:
                    if editable_pptx_error:
                        st.button("📊 PPT 생성 오류", disabled=True, use_container_width=True)
                        st.caption(f"PPT 생성 오류: {editable_pptx_error}")
                    else:
                        st.download_button(
                            label="📊 편집 가능한 PPT 다운로드",
                            data=editable_pptx_bytes,
                            file_name=f"{export_file_base}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                        )
                with col_d3:
                    st.button("🔄 입력 내용 전체 초기화", on_click=reset_all_fields, use_container_width=True)

        with tab_code:
            st.text_area("생성된 HTML 소스코드 (텍스트 소스 보관용)", value=final_mail_html, height=760)
